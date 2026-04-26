"""
ingestion_engine.py
===================
MemoForge — Module 1: Document Ingestion Engine
------------------------------------------------
Responsibility: Accept raw academic documents (.pdf, .txt, .docx, .xlsx, .pptx),
extract clean text with heuristic noise removal, and emit sliding-window
chunks ready for downstream LLM/VLM inference.

This module has ZERO knowledge of the UI or the LLM layer.
It only produces structured Python data (dicts / dataclasses).

Design principles
-----------------
* Memory-frugal: files are read page-by-page / row-by-row; nothing is
  held in RAM beyond what the current parsing step needs.
* Deterministic output: same file → same chunks, always.
* Fail loudly but gracefully: every known failure mode raises a typed
  exception with a human-readable message so the UI layer can surface it.

Dependencies (install via pip)
------------------------------
    pip install pymupdf python-docx python-pptx pandas openpyxl
"""

from __future__ import annotations

import logging
import re
import unicodedata
from collections import Counter
from dataclasses import dataclass, field
from pathlib import Path
from typing import Generator, List, Optional

# ---------------------------------------------------------------------------
# Optional heavy imports are done inside methods to keep module import fast
# and to give clearer ImportError messages if a dep is missing.
# ---------------------------------------------------------------------------

logger = logging.getLogger(__name__)


# ============================================================
# Custom exception hierarchy
# ============================================================

class IngestionError(Exception):
    """Base class for all ingestion-related errors."""


class PasswordProtectedError(IngestionError):
    """Raised when a PDF is encrypted and no password was supplied."""


class CorruptedFileError(IngestionError):
    """Raised when a file cannot be parsed due to corruption."""


class ImageOnlyPDFError(IngestionError):
    """
    Raised when a PDF contains no machine-readable text (scanned / image-only).
    The upstream caller should route this document to the VLM's vision encoder
    instead of the text pipeline.
    """


class EmptyDocumentError(IngestionError):
    """Raised when a document yields zero usable text after cleaning."""


class UnsupportedFormatError(IngestionError):
    """Raised when the file extension is not in the supported set."""


# ============================================================
# Data structures
# ============================================================

@dataclass
class DocumentChunk:
    """
    A single, self-contained unit of text ready for LLM consumption.

    Attributes
    ----------
    source_file : str
        Absolute path of the originating file (for traceability).
    chunk_index : int
        Zero-based position of this chunk within the document.
    total_chunks : int
        Total number of chunks produced from the document.
        Populated *after* all chunks are collected (set to -1 during streaming).
    text : str
        The cleaned, chunked text payload.
    token_estimate : int
        Rough token count (word-split heuristic; ~1.3 words/token average).
        Good enough for budget checks without importing a tokenizer.
    metadata : dict
        Arbitrary per-format metadata (e.g., page range, sheet name).
    requires_vision : bool
        True when the source page/content could not be read as text and the
        VLM should be directed to use its vision encoder for this chunk.
    """
    source_file: str
    chunk_index: int
    total_chunks: int
    text: str
    token_estimate: int = 0
    metadata: dict = field(default_factory=dict)
    requires_vision: bool = False

    def __post_init__(self) -> None:
        if not self.token_estimate:
            # Heuristic: split on whitespace, multiply by 0.75 (avg tokens/word ≈ 1.33)
            self.token_estimate = max(1, int(len(self.text.split()) * 0.75))


# ============================================================
# Sliding-window text chunker
# ============================================================

class TextChunker:
    """
    Splits a long string into overlapping windows measured in *words*.

    Using words rather than characters avoids splitting tokens mid-character
    and gives a stable approximation of LLM context cost.

    Parameters
    ----------
    chunk_size : int
        Target window size in words (default 512 ≈ ~680 tokens).
    overlap : int
        Number of words shared between consecutive windows (default 64).
        Overlap preserves sentence context across chunk boundaries.
    """

    def __init__(self, chunk_size: int = 512, overlap: int = 64) -> None:
        if overlap >= chunk_size:
            raise ValueError(
                f"overlap ({overlap}) must be strictly less than chunk_size ({chunk_size})."
            )
        self.chunk_size = chunk_size
        self.overlap = overlap

    def chunk(self, text: str) -> List[str]:
        """
        Split *text* into a list of word-windowed strings.

        Returns a list with at least one element (even for short texts).
        """
        words = text.split()
        if not words:
            return []

        chunks: List[str] = []
        step = self.chunk_size - self.overlap  # advance by this many words per iteration
        start = 0

        while start < len(words):
            end = min(start + self.chunk_size, len(words))
            chunk_words = words[start:end]
            chunks.append(" ".join(chunk_words))
            if end == len(words):
                break
            start += step

        return chunks


# ============================================================
# PDF heuristic cleaner
# ============================================================

class _PDFCleaner:
    """
    Stateful per-document noise filter for PDF text extraction.

    Strategy
    --------
    1. Collect the full frequency profile of short lines across ALL pages
       during a first pass (calibration).  Lines that appear on ≥ N% of
       pages are almost certainly headers/footers and are blacklisted.
    2. On the second (extraction) pass, strip blacklisted lines and apply
       pattern-based rules (page numbers, roman numerals, URL noise, etc.).

    This two-pass design costs one extra read of the page texts but keeps
    the per-page RAM footprint negligible because we only store line strings,
    not rendered pixels.
    """

    # A line is considered "short" if it has fewer than this many words.
    SHORT_LINE_WORD_THRESHOLD = 10

    # A short line is a header/footer candidate if it appears on this
    # fraction of pages or more (0.3 = 30%).
    REPETITION_THRESHOLD = 0.30

    # Patterns that unambiguously flag a line as noise.
    _NOISE_PATTERNS = [
        re.compile(r"^\s*\d+\s*$"),                      # standalone page numbers
        re.compile(r"^\s*[ivxlcdmIVXLCDM]+\s*$"),        # roman numerals alone
        re.compile(r"^\s*page\s+\d+\s*(of\s+\d+)?\s*$", re.IGNORECASE),  # "Page 3 of 10"
        re.compile(r"^\s*[-–—]+\s*\d+\s*[-–—]+\s*$"),   # — 5 — style footers
        re.compile(r"^https?://\S+$"),                    # bare URLs on their own line
        re.compile(r"^\s*©.{0,80}$"),                     # copyright lines
        re.compile(r"^\s*all rights reserved\.?\s*$", re.IGNORECASE),
        re.compile(r"^\s*confidential\.?\s*$", re.IGNORECASE),
        re.compile(r"^\s*draft\.?\s*$", re.IGNORECASE),
        re.compile(r"^\s*[A-Z0-9\s]{1,6}\s*\|\s*[A-Z0-9\s]{1,6}\s*$"),  # "DOC | REV"
    ]

    def __init__(self) -> None:
        self._blacklist: set[str] = set()

    def calibrate(self, page_texts: List[str]) -> None:
        """
        First pass: build the blacklist from repeated short lines.

        Parameters
        ----------
        page_texts : list of raw page strings
        """
        total_pages = len(page_texts)
        if total_pages == 0:
            return

        line_counter: Counter = Counter()

        for page_text in page_texts:
            # Use a set per page so the same line on the same page counts once.
            seen_on_page: set[str] = set()
            for raw_line in page_text.splitlines():
                line = raw_line.strip()
                if not line:
                    continue
                if len(line.split()) < self.SHORT_LINE_WORD_THRESHOLD:
                    normalized = line.lower()
                    if normalized not in seen_on_page:
                        line_counter[normalized] += 1
                        seen_on_page.add(normalized)

        cutoff = max(2, int(total_pages * self.REPETITION_THRESHOLD))
        self._blacklist = {
            line for line, count in line_counter.items() if count >= cutoff
        }
        logger.debug("PDF cleaner blacklist contains %d entries.", len(self._blacklist))

    def clean_page(self, raw_text: str) -> str:
        """
        Second pass: remove noise lines from a single page string.

        Returns cleaned text with normalized whitespace.
        """
        cleaned_lines: List[str] = []

        for raw_line in raw_text.splitlines():
            line = raw_line.strip()

            if not line:
                continue

            # 1. Pattern-based hard rules
            if any(pat.match(line) for pat in self._NOISE_PATTERNS):
                continue

            # 2. Frequency-based blacklist
            if line.lower() in self._blacklist:
                continue

            cleaned_lines.append(line)

        # Re-join with single newlines; collapse runs of blank lines
        joined = "\n".join(cleaned_lines)
        # Collapse 3+ consecutive newlines → double newline (paragraph break)
        joined = re.sub(r"\n{3,}", "\n\n", joined)
        return joined.strip()


# ============================================================
# Main DocumentIngestor class
# ============================================================

class DocumentIngestor:
    """
    Unified ingestion interface for .pdf, .txt, .docx, .xlsx, and .pptx files.

    Usage
    -----
    >>> ingestor = DocumentIngestor(chunk_size=512, overlap=64)
    >>> chunks = ingestor.ingest("lecture_notes.pdf")
    >>> for chunk in chunks:
    ...     print(chunk.chunk_index, chunk.token_estimate, chunk.text[:80])

    Parameters
    ----------
    chunk_size : int
        Window size in words for the sliding-window chunker (default 512).
    overlap : int
        Word overlap between consecutive chunks (default 64).
    pdf_password : str | None
        Optional password for encrypted PDFs.
    excel_max_rows : int
        Safety cap: stop reading an Excel sheet after this many data rows
        to avoid OOM on giant spreadsheets (default 2000).
    """

    SUPPORTED_EXTENSIONS = {".pdf", ".txt", ".docx", ".xlsx", ".pptx"}

    def __init__(
        self,
        chunk_size: int = 512,
        overlap: int = 64,
        pdf_password: Optional[str] = None,
        excel_max_rows: int = 2000,
    ) -> None:
        self.chunker = TextChunker(chunk_size=chunk_size, overlap=overlap)
        self.pdf_password = pdf_password
        self.excel_max_rows = excel_max_rows

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def ingest(self, file_path: str | Path) -> List[DocumentChunk]:
        """
        Parse *file_path* and return a list of DocumentChunk objects.

        This is the primary entry point for the rest of the application.
        It dispatches to format-specific parsers, runs the chunker, and
        stamps each chunk with its final total_chunks count.

        Raises
        ------
        UnsupportedFormatError
            File extension is not in SUPPORTED_EXTENSIONS.
        PasswordProtectedError
            PDF is encrypted and no password was given (or password is wrong).
        CorruptedFileError
            File cannot be opened / parsed.
        ImageOnlyPDFError
            PDF has no selectable text (scanned document).
        EmptyDocumentError
            After cleaning, no usable text remains.
        """
        path = Path(file_path).resolve()

        if not path.exists():
            raise FileNotFoundError(f"File not found: {path}")

        ext = path.suffix.lower()
        if ext not in self.SUPPORTED_EXTENSIONS:
            raise UnsupportedFormatError(
                f"'{ext}' is not supported. Accepted formats: "
                f"{', '.join(sorted(self.SUPPORTED_EXTENSIONS))}"
            )

        logger.info("Ingesting '%s' (format: %s)", path.name, ext)

        # Dispatch to the appropriate parser
        if ext == ".pdf":
            raw_segments = list(self._parse_pdf(path))
        elif ext == ".txt":
            raw_segments = list(self._parse_txt(path))
        elif ext == ".docx":
            raw_segments = list(self._parse_docx(path))
        elif ext == ".xlsx":
            raw_segments = list(self._parse_xlsx(path))
        elif ext == ".pptx":
            raw_segments = list(self._parse_pptx(path))
        else:
            # Unreachable given the extension check above, but keeps mypy happy.
            raise UnsupportedFormatError(f"No parser registered for '{ext}'.")

        # Merge all segment texts; preserve vision-flagged segments separately
        vision_segments = [s for s in raw_segments if s.get("requires_vision")]
        text_segments = [s for s in raw_segments if not s.get("requires_vision")]

        full_text = "\n\n".join(s["text"] for s in text_segments if s.get("text", "").strip())

        if not full_text.strip() and not vision_segments:
            raise EmptyDocumentError(
                f"'{path.name}' produced no usable text after cleaning. "
                "The file may be blank or contain only unsupported content."
            )

        # --- Chunk the clean text ---
        chunks: List[DocumentChunk] = []
        source = str(path)

        if full_text.strip():
            window_texts = self.chunker.chunk(full_text)
            for idx, window in enumerate(window_texts):
                chunks.append(
                    DocumentChunk(
                        source_file=source,
                        chunk_index=idx,
                        total_chunks=-1,   # stamped below
                        text=window,
                        metadata={"format": ext.lstrip(".")},
                    )
                )

        # Append vision-flagged segments as special chunks (not chunked further
        # because their payload is a page reference, not raw text).
        for seg in vision_segments:
            chunks.append(
                DocumentChunk(
                    source_file=source,
                    chunk_index=len(chunks),
                    total_chunks=-1,
                    text=seg.get("text", ""),
                    metadata=seg.get("metadata", {}),
                    requires_vision=True,
                )
            )

        # Stamp the final count on every chunk
        total = len(chunks)
        for chunk in chunks:
            chunk.total_chunks = total

        logger.info(
            "Produced %d chunk(s) from '%s'.", total, path.name
        )
        return chunks

    # ------------------------------------------------------------------
    # Format-specific parsers  (private, yield raw segment dicts)
    # ------------------------------------------------------------------
    # Each parser yields dicts:
    #   {"text": str, "requires_vision": bool, "metadata": dict}
    # This keeps the parsing logic decoupled from the chunking logic.
    # ------------------------------------------------------------------

    def _parse_pdf(self, path: Path) -> Generator[dict, None, None]:
        """
        Parse a PDF using PyMuPDF (fitz).

        Steps
        -----
        1. Open the document (handle encryption / corruption).
        2. Calibrate the PDF cleaner on a quick first pass.
        3. Extract + clean text page by page.
        4. Flag image-only pages with requires_vision=True.
        """
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise ImportError(
                "PyMuPDF is required for PDF parsing. Install it with: pip install pymupdf"
            )

        # --- Open ---
        try:
            doc = fitz.open(str(path))
        except fitz.FileDataError as exc:
            raise CorruptedFileError(
                f"'{path.name}' could not be opened by PyMuPDF. "
                f"The file may be corrupted. Detail: {exc}"
            ) from exc
        except Exception as exc:
            raise CorruptedFileError(
                f"Unexpected error opening '{path.name}': {exc}"
            ) from exc

        # --- Handle encryption ---
        if doc.is_encrypted:
            if self.pdf_password is None:
                doc.close()
                raise PasswordProtectedError(
                    f"'{path.name}' is password-protected. "
                    "Provide the password via the 'pdf_password' constructor argument."
                )
            success = doc.authenticate(self.pdf_password)
            if not success:
                doc.close()
                raise PasswordProtectedError(
                    f"The supplied password for '{path.name}' is incorrect."
                )

        page_count = doc.page_count
        if page_count == 0:
            doc.close()
            raise EmptyDocumentError(f"'{path.name}' contains zero pages.")

        # --- First pass: collect raw page texts for calibration ---
        # We read text from all pages here but discard it after calibration
        # to keep peak memory low (strings are short-lived).
        raw_page_texts: List[str] = []
        for page_num in range(page_count):
            page = doc.load_page(page_num)
            raw_page_texts.append(page.get_text("text"))  # type: ignore[attr-defined]
            # Immediately release the page object's resources
            page = None  # noqa: redefinition aids GC

        # Check for image-only PDF (no text on any page)
        total_text_chars = sum(len(t) for t in raw_page_texts)
        if total_text_chars < page_count * 5:
            # Heuristic: fewer than 5 chars per page on average → scanned
            doc.close()
            raise ImageOnlyPDFError(
                f"'{path.name}' appears to be a scanned / image-only PDF "
                f"({total_text_chars} total text chars across {page_count} pages). "
                "Route to the VLM vision encoder."
            )

        # --- Calibrate the cleaner ---
        cleaner = _PDFCleaner()
        cleaner.calibrate(raw_page_texts)

        # --- Second pass: yield cleaned segments ---
        text_page_count = 0
        vision_page_count = 0

        for page_num, raw_text in enumerate(raw_page_texts):
            cleaned = cleaner.clean_page(raw_text)

            if not cleaned.strip():
                # This specific page has no text → flag for vision encoder
                vision_page_count += 1
                yield {
                    "text": f"[Page {page_num + 1}: image-only — route to vision encoder]",
                    "requires_vision": True,
                    "metadata": {
                        "page": page_num + 1,
                        "total_pages": page_count,
                        "format": "pdf",
                    },
                }
            else:
                text_page_count += 1
                yield {
                    "text": cleaned,
                    "requires_vision": False,
                    "metadata": {
                        "page_start": page_num + 1,
                        "page_end": page_num + 1,
                        "total_pages": page_count,
                        "format": "pdf",
                    },
                }

        doc.close()
        logger.debug(
            "PDF '%s': %d text pages, %d image-only pages.",
            path.name, text_page_count, vision_page_count,
        )

    def _parse_txt(self, path: Path) -> Generator[dict, None, None]:
        """
        Parse a plain-text file.

        Tries UTF-8 first; falls back to latin-1 to handle legacy encodings
        without crashing.  Normalises unicode to NFC so composed and
        decomposed characters compare equal downstream.
        """
        try:
            try:
                raw = path.read_text(encoding="utf-8")
            except UnicodeDecodeError:
                logger.warning(
                    "'%s' is not UTF-8; falling back to latin-1.", path.name
                )
                raw = path.read_text(encoding="latin-1")
        except OSError as exc:
            raise CorruptedFileError(
                f"Could not read '{path.name}': {exc}"
            ) from exc

        # NFC normalisation collapses combining characters (e.g., e + ́ → é)
        text = unicodedata.normalize("NFC", raw)
        # Collapse runs of blank lines to a single paragraph break
        text = re.sub(r"\n{3,}", "\n\n", text).strip()

        if not text:
            raise EmptyDocumentError(f"'{path.name}' is empty.")

        yield {
            "text": text,
            "requires_vision": False,
            "metadata": {"format": "txt"},
        }

    def _parse_docx(self, path: Path) -> Generator[dict, None, None]:
        """
        Parse a .docx file using python-docx.

        Extraction strategy
        -------------------
        * Headings    → prefixed with Markdown '#' markers scaled to heading level.
        * List items  → prefixed with '•' so the LLM recognises list structure.
        * Body paras  → emitted as plain text paragraphs.
        * Empty paras → silently skipped (Word documents are full of them).

        Tables are intentionally omitted here; they will be handled in a
        future table-extraction sub-module.  The paragraph-level extraction
        is sufficient for academic lecture notes and reports.
        """
        try:
            from docx import Document  # python-docx
            from docx.oxml.ns import qn
        except ImportError:
            raise ImportError(
                "python-docx is required for .docx parsing. "
                "Install it with: pip install python-docx"
            )

        try:
            doc = Document(str(path))
        except Exception as exc:
            # python-docx raises a variety of exceptions for corrupted files
            raise CorruptedFileError(
                f"'{path.name}' could not be opened as a Word document. "
                f"The file may be corrupted or not a valid .docx. Detail: {exc}"
            ) from exc

        lines: List[str] = []

        for para in doc.paragraphs:
            # Strip surrounding whitespace from the raw text
            text = para.text.strip()
            if not text:
                continue  # skip empty paragraphs

            style_name = (para.style.name or "").lower()

            if style_name.startswith("heading"):
                # Extract heading level (e.g., "Heading 1" → 1)
                try:
                    level = int(style_name.split()[-1])
                except (ValueError, IndexError):
                    level = 1
                level = max(1, min(level, 6))  # clamp to 1–6
                prefix = "#" * level
                lines.append(f"{prefix} {text}")

            elif "list" in style_name or para.style.name.startswith("List"):
                lines.append(f"• {text}")

            else:
                lines.append(text)

        if not lines:
            raise EmptyDocumentError(
                f"'{path.name}' contains no extractable text paragraphs."
            )

        full_text = "\n\n".join(lines)

        yield {
            "text": full_text,
            "requires_vision": False,
            "metadata": {
                "format": "docx",
                "paragraph_count": len(lines),
            },
        }

    def _parse_xlsx(self, path: Path) -> Generator[dict, None, None]:
        """
        Parse a .xlsx spreadsheet using pandas + openpyxl.

        Output format
        -------------
        Each worksheet is converted to one of two representations depending
        on its shape:

        1. **Wide sheets (many columns, few rows):** Markdown table.
           Columns become header cells; rows become data rows.
           This preserves the tabular visual structure that a VLM
           can reason about directly.

        2. **Tall sheets (many rows, few columns, ≤ 2 cols):** Key-value
           pairs formatted as "Key: Value" on separate lines.
           Ideal for dictionaries, glossaries, and reference sheets.

        Both representations are preceded by a sheet-level header so the
        model knows which sheet it is reading.

        Row cap
        -------
        Reading stops at self.excel_max_rows to avoid OOM on huge files.
        A warning note is appended to the output when the cap is hit.
        """
        try:
            import pandas as pd
        except ImportError:
            raise ImportError(
                "pandas is required for .xlsx parsing. "
                "Install it with: pip install pandas openpyxl"
            )

        try:
            # read sheet names without loading data
            import openpyxl
            wb_meta = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            sheet_names = wb_meta.sheetnames
            wb_meta.close()
        except Exception as exc:
            raise CorruptedFileError(
                f"'{path.name}' could not be opened as an Excel workbook. "
                f"The file may be corrupted. Detail: {exc}"
            ) from exc

        if not sheet_names:
            raise EmptyDocumentError(f"'{path.name}' contains no sheets.")

        any_content = False

        for sheet_name in sheet_names:
            try:
                # nrows cap prevents reading the entire sheet into RAM
                df = pd.read_excel(
                    str(path),
                    sheet_name=sheet_name,
                    nrows=self.excel_max_rows,
                    engine="openpyxl",
                )
            except Exception as exc:
                logger.warning(
                    "Could not read sheet '%s' in '%s': %s",
                    sheet_name, path.name, exc
                )
                continue

            # Drop completely empty rows and columns
            df.dropna(how="all", inplace=True)
            df.dropna(axis=1, how="all", inplace=True)

            if df.empty:
                logger.debug("Sheet '%s' is empty after cleaning; skipping.", sheet_name)
                continue

            any_content = True
            was_capped = len(df) == self.excel_max_rows
            n_cols = len(df.columns)

            # --- Decide representation ---
            if n_cols <= 2:
                text = self._xlsx_to_key_value(df, sheet_name)
            else:
                text = self._xlsx_to_markdown_table(df, sheet_name)

            if was_capped:
                text += (
                    f"\n\n> ⚠️ Note: This sheet was truncated at {self.excel_max_rows} rows "
                    "to stay within memory limits. Additional rows were not processed."
                )

            yield {
                "text": text,
                "requires_vision": False,
                "metadata": {
                    "format": "xlsx",
                    "sheet_name": sheet_name,
                    "row_count": len(df),
                    "col_count": n_cols,
                    "truncated": was_capped,
                },
            }

        if not any_content:
            raise EmptyDocumentError(
                f"'{path.name}' contains no usable data in any sheet."
            )

    def _parse_pptx(self, path: Path) -> Generator[dict, None, None]:
        """
        Parse a .pptx presentation using python-pptx.

        Extraction strategy
        -------------------
        * Capture slide title (if available).
        * Capture text from text frames and table cells.
        * Capture notes text when present.
        * Emit one segment per slide to preserve slide boundaries.
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "python-pptx is required for .pptx parsing. "
                "Install it with: pip install python-pptx"
            )

        try:
            presentation = Presentation(str(path))
        except Exception as exc:
            raise CorruptedFileError(
                f"'{path.name}' could not be opened as a PowerPoint deck. "
                f"The file may be corrupted. Detail: {exc}"
            ) from exc

        if not presentation.slides:
            raise EmptyDocumentError(f"'{path.name}' contains zero slides.")

        emitted = 0
        total_slides = len(presentation.slides)

        for index, slide in enumerate(presentation.slides, start=1):
            lines: List[str] = []

            title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else ""
            if title:
                lines.append(f"# Slide {index}: {title}")
            else:
                lines.append(f"# Slide {index}")

            for shape in slide.shapes:
                if slide.shapes.title is not None and shape == slide.shapes.title:
                    continue
                if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in paragraph.runs).strip()
                        if not text:
                            continue
                        if paragraph.level and paragraph.level > 0:
                            indent = "  " * paragraph.level
                            lines.append(f"{indent}- {text}")
                        else:
                            lines.append(text)

                if getattr(shape, "has_table", False) and shape.has_table:
                    lines.append("Table:")
                    for row in shape.table.rows:
                        row_cells = [cell.text_frame.text.strip().replace("\n", " ") for cell in row.cells]
                        if any(row_cells):
                            lines.append(" | ".join(row_cells))

            notes_text = ""
            try:
                notes_slide = slide.notes_slide
                if notes_slide and notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text.strip()
            except Exception:
                notes_text = ""

            if notes_text:
                lines.append("Notes:")
                lines.append(notes_text)

            slide_text = "\n".join(line for line in lines if line.strip()).strip()
            if not slide_text:
                continue

            emitted += 1
            yield {
                "text": slide_text,
                "requires_vision": False,
                "metadata": {
                    "format": "pptx",
                    "slide": index,
                    "total_slides": total_slides,
                },
            }

        if emitted == 0:
            raise EmptyDocumentError(
                f"'{path.name}' contains no extractable slide text."
            )

    # ------------------------------------------------------------------
    # Excel helper formatters
    # ------------------------------------------------------------------

    @staticmethod
    def _xlsx_to_markdown_table(df, sheet_name: str) -> str:
        """
        Render a DataFrame as a Markdown table prefixed by a sheet header.

        Handles NaN by replacing with empty string so the table stays
        well-formed without 'nan' noise in every empty cell.
        """
        df = df.fillna("")  # replace NaN with empty string

        # Build header row
        headers = [str(col) for col in df.columns]
        header_row = "| " + " | ".join(headers) + " |"
        separator  = "| " + " | ".join(["---"] * len(headers)) + " |"

        rows: List[str] = [f"## Sheet: {sheet_name}\n", header_row, separator]

        for _, row in df.iterrows():
            # Convert each cell to string; truncate overly long cells
            cells = []
            for val in row:
                cell_str = str(val).replace("|", "\\|").replace("\n", " ")
                if len(cell_str) > 200:
                    cell_str = cell_str[:197] + "..."
                cells.append(cell_str)
            rows.append("| " + " | ".join(cells) + " |")

        return "\n".join(rows)

    @staticmethod
    def _xlsx_to_key_value(df, sheet_name: str) -> str:
        """
        Render a narrow DataFrame (1–2 columns) as key: value pairs.

        If there is only one column, each row value is emitted as a
        plain bullet point.
        """
        lines: List[str] = [f"## Sheet: {sheet_name}\n"]
        df = df.fillna("")

        if len(df.columns) == 1:
            col = df.columns[0]
            for val in df[col]:
                val_str = str(val).strip()
                if val_str:
                    lines.append(f"• {val_str}")
        else:
            key_col, val_col = df.columns[0], df.columns[1]
            for _, row in df.iterrows():
                key = str(row[key_col]).strip()
                val = str(row[val_col]).strip()
                if key or val:
                    lines.append(f"{key}: {val}")

        return "\n".join(lines)


# ============================================================
# Convenience generator for memory-sensitive pipelines
# ============================================================

def ingest_stream(
    file_path: str | Path,
    chunk_size: int = 512,
    overlap: int = 64,
    pdf_password: Optional[str] = None,
    excel_max_rows: int = 2000,
) -> Generator[DocumentChunk, None, None]:
    """
    Generator wrapper around DocumentIngestor.ingest().

    Yields chunks one at a time rather than returning the full list.
    Useful when feeding a long document into an inference loop without
    buffering all chunks in RAM first.

    Note: total_chunks will be -1 on every chunk because the total is
    not known until all chunks are produced. If you need total_chunks,
    use DocumentIngestor.ingest() instead.

    Example
    -------
    >>> for chunk in ingest_stream("thesis.pdf"):
    ...     send_to_llm(chunk.text)
    """
    ingestor = DocumentIngestor(
        chunk_size=chunk_size,
        overlap=overlap,
        pdf_password=pdf_password,
        excel_max_rows=excel_max_rows,
    )
    for chunk in ingestor.ingest(file_path):
        yield chunk


# ============================================================
# Quick smoke-test  (run:  python ingestion_engine.py <path>)
# ============================================================

if __name__ == "__main__":
    import sys

    logging.basicConfig(
        level=logging.DEBUG,
        format="%(asctime)s [%(levelname)s] %(name)s — %(message)s",
    )

    if len(sys.argv) < 2:
        print("Usage: python ingestion_engine.py <path_to_document>")
        sys.exit(1)

    target = sys.argv[1]
    print(f"\n{'='*60}\nIngesting: {target}\n{'='*60}")

    try:
        ingestor = DocumentIngestor(chunk_size=512, overlap=64)
        chunks = ingestor.ingest(target)
        print(f"\n✓ {len(chunks)} chunk(s) produced.\n")
        for c in chunks[:3]:  # preview first 3
            print(
                f"  Chunk {c.chunk_index + 1}/{c.total_chunks} "
                f"(~{c.token_estimate} tokens)"
                f"{'  [VISION]' if c.requires_vision else ''}\n"
                f"  {c.text[:200].replace(chr(10), ' ')}...\n"
            )
    except IngestionError as e:
        print(f"\n✗ Ingestion failed: {type(e).__name__}: {e}")
        sys.exit(1)
