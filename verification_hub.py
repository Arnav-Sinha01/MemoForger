"""
verification_hub.py
===================
MemoForge â€” Module 4: Verification & Export Hub
-------------------------------------------------
Responsibility: Present the structured output from Module 2 (VLMPipeline) in an
interactive review window. The user can browse, filter, inline-edit, and delete
flashcards before exporting them in multiple formats.

This module has ZERO knowledge of the pipeline or ingestion layers.
It accepts a List[GenerationResult] and returns nothing â€” all side-effects
(file saves, clipboard writes) are confined to this module's export helpers.

Design Language: "Academic Brutalism" (identical to Module 3)
--------------------------------------------------------------
All _Theme constants are redeclared here verbatim so Module 4 can be
developed / tested independently.  If the palette ever changes, update
_Theme in app_frame.py AND here (a single grep + replace operation).

Layout (annotated wireframe)
-----------------------------
â”Œâ”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”
â”‚  â–‘ MEMOFORGE â€” VERIFICATION & EXPORT HUB            [âœ• CLOSE]   â”‚  â† TitleBar
â”œâ”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€-â”¤
â”‚  STATS BANNER:  13 cards  Â·  2 sources  Â·  1 chunk failed        â”‚
â”œâ”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”¬â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€-â”¤
â”‚                    â”‚                                              â”‚
â”‚  FILTER PANEL      â”‚  FLASHCARD GRID (scrollable, 2-col)         â”‚
â”‚  â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€     â”‚  â”Œâ”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”  â”Œâ”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”  â”‚
â”‚  SOURCES           â”‚  â”‚ [EASY]          â”‚  â”‚ [HARD]           â”‚  â”‚
â”‚  â˜‘ doc.pdf  (8)    â”‚  â”‚ Q: ...          â”‚  â”‚ Q: ...           â”‚  â”‚
â”‚  â˜‘ notes.txt (5)   â”‚  â”‚ A: ...          â”‚  â”‚ A: ...           â”‚  â”‚
â”‚                    â”‚  â”‚ â—‹bio â—‹cell      â”‚  â”‚ â—‹physics â—‹force  â”‚  â”‚
â”‚  DIFFICULTY        â”‚  â”‚ [EDIT]    [âœ•]   â”‚  â”‚ [EDIT]    [âœ•]    â”‚  â”‚
â”‚  â˜‘ easy   (3)      â”‚  â””â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”˜  â””â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”˜  â”‚
â”‚  â˜‘ medium (7)      â”‚                                              â”‚
â”‚  â˜‘ hard   (3)      â”‚  (failed results shown as collapsed banners) â”‚
â”‚                    â”‚                                              â”‚
â”‚  TAGS              â”‚                                              â”‚
â”‚  â˜‘ biology         â”‚                                              â”‚
â”‚  â˜‘ cell-theory     â”‚                                              â”‚
â”‚                    â”‚                                              â”‚
â”œâ”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”´â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€-â”¤
â”‚  [FORMAT â–¼] [â¬‡ CONFIRM & SAVE]          13 cards Â· 0 deleted    â”‚  â† ExportBar
â””â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”˜

Export Formats (via _ExportEngine)
------------------------------------
  DOCX (.docx), PDF (.pdf), and PowerPoint (.pptx)
  Anki Deck (.apkg), CSV (.csv), and JSON (.json)

Integration Point (called by Module 5)
---------------------------------------
    from verification_hub import VerificationHub

    hub = VerificationHub(parent=app, results=generation_results)
    hub.grab_set()   # modal
    hub.wait_window()

Dependencies
------------
    pip install customtkinter genanki
"""

from __future__ import annotations

import csv
import json
import logging
import random
import re
import tkinter as tk
from datetime import datetime
from dataclasses import dataclass, field
from pathlib import Path
from tkinter import filedialog, messagebox
from typing import Callable, Dict, List, Optional, Set

import customtkinter as ctk

# ---------------------------------------------------------------------------
# genanki â€” graceful degradation if the library is not installed.
# The rest of the module (CSV / JSON export, UI) works without it.
# Only Anki export raises ImportError at runtime if genanki is missing.
# ---------------------------------------------------------------------------
try:
    import genanki  # type: ignore
    _GENANKI_AVAILABLE = True
except ImportError:
    genanki = None          # type: ignore
    _GENANKI_AVAILABLE = False

logger = logging.getLogger(__name__)


# ============================================================
# Graceful import of Module 2 data structures
# ============================================================

try:
    from vlm_pipeline import Flashcard, GenerationResult  # type: ignore
except ImportError:
    # Stub types so Module 4 can be run / tested standalone.
    from dataclasses import dataclass as _dc, field as _f
    from typing import List as _L

    @_dc
    class Flashcard:  # type: ignore[no-redef]
        question: str
        answer: str
        difficulty: str = "medium"
        tags: _L[str] = _f(default_factory=list)
        source_chunk_index: int = -1

        def to_dict(self) -> dict:
            return {
                "question":           self.question,
                "answer":             self.answer,
                "difficulty":         self.difficulty,
                "tags":               self.tags,
                "source_chunk_index": self.source_chunk_index,
            }

    @_dc
    class GenerationResult:  # type: ignore[no-redef]
        chunk_index: int
        source_file: str
        flashcards: _L[Flashcard] = _f(default_factory=list)
        summary: str = ""
        was_vision_chunk: bool = False
        error: Optional[str] = None
        latency_seconds: float = 0.0

        @property
        def succeeded(self) -> bool:
            return self.error is None

    logger.warning("vlm_pipeline not found â€” running Module 4 with stub types.")


# ============================================================
# Design tokens (mirror of app_frame._Theme)
# ============================================================

class _Theme:
    # Colour palette (dark mode)
    BG_ROOT         = "#0d1117"
    BG_PANEL        = "#161b22"
    BG_CARD         = "#21262d"
    BG_DROPZONE     = "#0d1117"
    BG_DROPZONE_HV  = "#1f2428"
    BG_EDIT         = "#0d1117"   # card edit-mode tint

    ACCENT          = "#2f81f7"
    ACCENT_HOVER    = "#58a6ff"
    ACCENT_PRESS    = "#1f6feb"

    TEXT_PRIMARY    = "#e6edf3"
    TEXT_SECONDARY  = "#8b949e"
    TEXT_MONO       = "#f0f6fc"

    SUCCESS         = "#2ea043"
    ERROR           = "#f85149"
    WARNING         = "#d29922"

    BORDER          = "#30363d"
    BORDER_FOCUS    = "#58a6ff"

    # Typography
    FONT_DISPLAY    = ("Segoe UI", 24, "bold")
    FONT_HEADING    = ("Segoe UI", 12, "bold")
    FONT_SUBHEAD    = ("Segoe UI", 12, "bold")
    FONT_BODY       = ("Segoe UI", 12)
    FONT_SMALL      = ("Segoe UI", 11)
    FONT_MONO       = ("Consolas", 11)
    FONT_BTN        = ("Segoe UI", 13, "bold")
    FONT_CARD_Q     = ("Segoe UI", 13, "bold")
    FONT_CARD_A     = ("Segoe UI", 12)

    # Geometry
    CORNER_RADIUS   = 8
    PAD_OUTER       = 24
    PAD_INNER       = 16
    PAD_SMALL       = 8
    BTN_HEIGHT      = 44
    FILTER_W        = 220
    CARD_W          = 320
    CARD_GAP        = 12

    # Difficulty colours
    DIFF_COLORS = {
        "easy":   "#2ea043",
        "medium": "#d29922",
        "hard":   "#f85149",
    }


# ============================================================
# Export format registry
# ============================================================
# Each entry is a dict describing one export format.  Adding a new
# format only requires adding an entry here and a matching method in
# _ExportEngine â€” the UI (dropdown, dialog) auto-derives from this list.

_EXPORT_FORMATS: List[Dict] = [
    {
        "label":     "DOCX (.docx)",
        "ext":       ".docx",
        "filetypes": [
            ("Word Document", "*.docx"),
            ("All files",     "*.*"),
        ],
    },
    {
        "label":     "PDF (.pdf)",
        "ext":       ".pdf",
        "filetypes": [
            ("PDF Document", "*.pdf"),
            ("All files",    "*.*"),
        ],
    },
    {
        "label":     "PowerPoint (.pptx)",
        "ext":       ".pptx",
        "filetypes": [
            ("PowerPoint", "*.pptx"),
            ("All files",  "*.*"),
        ],
    },
    {
        "label":     "Anki Deck (.apkg)",   # shown in the CTkOptionMenu
        "ext":       ".apkg",               # default extension for filedialog
        "filetypes": [                       # filter list for filedialog
            ("Anki Package", "*.apkg"),
            ("All files",    "*.*"),
        ],
    },
    {
        "label":     "CSV (.csv)",
        "ext":       ".csv",
        "filetypes": [
            ("CSV file",  "*.csv"),
            ("All files", "*.*"),
        ],
    },
    {
        "label":     "JSON (.json)",
        "ext":       ".json",
        "filetypes": [
            ("JSON file", "*.json"),
            ("All files", "*.*"),
        ],
    },
]

# Convenience structures derived from the registry above
_FORMAT_BY_LABEL: Dict[str, Dict] = {f["label"]: f for f in _EXPORT_FORMATS}
_FORMAT_LABELS:   List[str]       = [f["label"] for f in _EXPORT_FORMATS]


def _sanitize_llm_text(text: str) -> str:
    """Remove common model-formatting artifacts from visible/exported text."""
    cleaned = (text or "").strip()
    if not cleaned:
        return ""

    candidate = cleaned.strip()
    if candidate.startswith("{") and candidate.endswith("}"):
        try:
            parsed = json.loads(candidate)
            if isinstance(parsed, dict):
                for key in ("question", "answer", "summary", "front", "back", "text", "content", "note"):
                    value = parsed.get(key)
                    if isinstance(value, str) and value.strip():
                        cleaned = value.strip()
                        break
            elif isinstance(parsed, list) and parsed and isinstance(parsed[0], dict):
                first = parsed[0]
                for key in ("question", "answer", "summary", "front", "back", "text", "content", "note"):
                    value = first.get(key)
                    if isinstance(value, str) and value.strip():
                        cleaned = value.strip()
                        break
        except Exception:
            pass

    cleaned = re.sub(r"```(?:json)?|```", " ", cleaned, flags=re.IGNORECASE)
    cleaned = cleaned.replace("\\n", " ").replace("\\t", " ")
    cleaned = re.sub(r"^\s*(?:question|answer|summary|q|a)\s*[:\-]\s*", "", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"^\s*(?:json|output)\s*[:\-]\s*", "", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(
        r"\{?\s*\"?(?:question|answer|summary|flashcards|difficulty|tags|front|back|note|output|json)\"?\s*:\s*\"?[^{}]{0,260}\"?\s*\}?",
        " ",
        cleaned,
        flags=re.IGNORECASE,
    )
    cleaned = re.sub(r"\"[A-Za-z_][A-Za-z0-9_\-\s]{0,30}\"?\s*:\s*", "", cleaned)
    cleaned = re.sub(r"\b(?:question|answer|summary|front|back)\b\s*:\s*", "", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\s+", " ", cleaned).strip()
    cleaned = cleaned.strip("`").strip().strip('"').strip("'").strip()
    cleaned = re.sub(r"^[\[\{]+", "", cleaned)
    cleaned = re.sub(r"[\]\}]+$", "", cleaned)
    cleaned = cleaned.replace("{", " ").replace("}", " ")
    cleaned = re.sub(r"\s+", " ", cleaned).strip()
    return cleaned


def _sanitize_tag_text(tag: str) -> str:
    cleaned = _sanitize_llm_text(tag).lstrip("#")
    cleaned = re.sub(r"[^\w\s\-]", "", cleaned, flags=re.UNICODE)
    cleaned = re.sub(r"\s+", " ", cleaned).strip().lower()
    return cleaned[:40]


# ============================================================
# Internal data model
# ============================================================

@dataclass
class _CardRecord:
    """
    Mutable wrapper around a Flashcard that tracks edit/delete state.
    The verification grid works entirely with these records, never the
    raw Flashcard objects directly, so mutations never corrupt the source.
    """
    card: Flashcard
    source_file: str
    chunk_index: int
    deleted: bool = False

    # Live-editable copies (diverge from card after user edits)
    question:   str       = field(init=False)
    answer:     str       = field(init=False)
    tags:       List[str] = field(init=False)
    difficulty: str       = field(init=False)

    def __post_init__(self) -> None:
        self.question   = _sanitize_llm_text(self.card.question)
        self.answer     = _sanitize_llm_text(self.card.answer)
        self.tags       = [_sanitize_tag_text(tag) for tag in list(self.card.tags)]
        self.tags       = [tag for tag in self.tags if tag]
        self.difficulty = self.card.difficulty or "medium"

    def to_flashcard(self) -> Flashcard:
        """Return a new Flashcard reflecting any edits."""
        return Flashcard(
            question           = self.question,
            answer             = self.answer,
            difficulty         = self.difficulty,
            tags               = self.tags,
            source_chunk_index = self.chunk_index,
        )

    @property
    def source_name(self) -> str:
        return Path(self.source_file).name


@dataclass
class _SummaryRecord:
    source_file: str
    chunk_index: int
    summary: str

    @property
    def source_name(self) -> str:
        return Path(self.source_file).name


# ============================================================
# Stats banner
# ============================================================

class _StatsBanner(ctk.CTkFrame):
    """
    Full-width strip at the top of the hub showing aggregate numbers.
    Refreshed by VerificationHub whenever the card set changes.
    """

    def __init__(self, parent, **kwargs) -> None:
        super().__init__(
            parent,
            fg_color      = _Theme.BG_PANEL,
            corner_radius = 0,
            **kwargs,
        )
        self._build()

    def _build(self) -> None:
        self.grid_columnconfigure(0, weight=1)

        sep = ctk.CTkFrame(self, fg_color=_Theme.BORDER, height=1, corner_radius=0)
        sep.grid(row=0, column=0, sticky="ew")

        inner = ctk.CTkFrame(self, fg_color="transparent")
        inner.grid(row=1, column=0, sticky="ew",
                   padx=_Theme.PAD_OUTER, pady=(_Theme.PAD_SMALL, _Theme.PAD_SMALL))
        inner.grid_columnconfigure(0, weight=1)

        self._stat_label = ctk.CTkLabel(
            inner,
            text       = "",
            font       = _Theme.FONT_MONO,
            text_color = _Theme.TEXT_SECONDARY,
            anchor     = "w",
        )
        self._stat_label.grid(row=0, column=0, sticky="w")

        sep2 = ctk.CTkFrame(self, fg_color=_Theme.BORDER, height=1, corner_radius=0)
        sep2.grid(row=2, column=0, sticky="ew")

    def refresh(
        self,
        total: int,
        visible: int,
        deleted: int,
        sources: int,
        failed_chunks: int,
        total_latency: float,
    ) -> None:
        parts = [
            f"{total} card(s) generated",
            f"{visible} visible",
            f"{deleted} deleted",
            f"{sources} source(s)",
        ]
        if failed_chunks:
            parts.append(f"{failed_chunks} chunk(s) failed")
        if total_latency > 0:
            parts.append(f"{total_latency:.1f}s inference")
        self._stat_label.configure(text=" | ".join(parts))


# ============================================================
# Filter sidebar
# ============================================================

class _FilterPanel(ctk.CTkScrollableFrame):
    """
    Left-hand sidebar with checkbox groups for source file, difficulty,
    and tag filtering.  Calls on_filter_change() on every toggle.
    """

    _DIFFICULTIES = ("easy", "medium", "hard")

    def __init__(
        self,
        parent,
        on_filter_change: Callable[[], None],
        **kwargs,
    ) -> None:
        super().__init__(
            parent,
            fg_color      = _Theme.BG_PANEL,
            corner_radius = 0,
            width         = _Theme.FILTER_W,
            scrollbar_button_color       = _Theme.BORDER,
            scrollbar_button_hover_color = _Theme.ACCENT,
            **kwargs,
        )

        self._on_change  = on_filter_change
        self._source_vars: Dict[str, tk.BooleanVar] = {}
        self._diff_vars:   Dict[str, tk.BooleanVar] = {}
        self._tag_vars:    Dict[str, tk.BooleanVar] = {}
        self._source_counts: Dict[str, int] = {}
        self._diff_counts:   Dict[str, int] = {}
        self._tag_counts:    Dict[str, int] = {}

    def populate(self, records: List[_CardRecord]) -> None:
        """(Re)build checkboxes from the full record list."""
        self._source_counts = {}
        self._diff_counts   = {d: 0 for d in self._DIFFICULTIES}
        self._tag_counts    = {}

        for r in records:
            self._source_counts[r.source_name] = (
                self._source_counts.get(r.source_name, 0) + 1
            )
            diff = r.difficulty.lower()
            if diff in self._diff_counts:
                self._diff_counts[diff] += 1
            for tag in r.tags:
                self._tag_counts[tag] = self._tag_counts.get(tag, 0) + 1

        for w in self.winfo_children():
            w.destroy()

        row = 0
        row = self._section("SOURCES", row)
        row = self._checkbox_group(self._source_counts, self._source_vars, row)
        row = self._section("DIFFICULTY", row)
        row = self._checkbox_group(
            {d: self._diff_counts[d] for d in self._DIFFICULTIES if self._diff_counts[d]},
            self._diff_vars, row, color_map=_Theme.DIFF_COLORS,
        )
        if self._tag_counts:
            row = self._section("TAGS", row)
            row = self._checkbox_group(self._tag_counts, self._tag_vars, row)

    @property
    def active_sources(self) -> Set[str]:
        return {k for k, v in self._source_vars.items() if v.get()}

    @property
    def active_difficulties(self) -> Set[str]:
        return {k for k, v in self._diff_vars.items() if v.get()}

    @property
    def active_tags(self) -> Set[str]:
        return {k for k, v in self._tag_vars.items() if v.get()}

    def select_all(self) -> None:
        for v in (*self._source_vars.values(),
                  *self._diff_vars.values(),
                  *self._tag_vars.values()):
            v.set(True)
        self._on_change()

    def _section(self, label: str, row: int) -> int:
        if row > 0:
            sp = ctk.CTkFrame(self, fg_color="transparent", height=8, corner_radius=0)
            sp.grid(row=row, column=0, sticky="ew")
            row += 1

        hdr = ctk.CTkLabel(
            self, text=label, font=_Theme.FONT_HEADING,
            text_color=_Theme.ACCENT, anchor="w",
        )
        hdr.grid(row=row, column=0, sticky="ew",
                 padx=(_Theme.PAD_INNER, 0), pady=(4, 2))
        row += 1

        sep = ctk.CTkFrame(self, fg_color=_Theme.BORDER, height=1, corner_radius=0)
        sep.grid(row=row, column=0, sticky="ew", padx=_Theme.PAD_INNER)
        row += 1
        return row

    def _checkbox_group(
        self,
        counts:    Dict[str, int],
        var_store: Dict[str, tk.BooleanVar],
        row:       int,
        color_map: Optional[Dict[str, str]] = None,
    ) -> int:
        existing_keys = set(var_store.keys())
        for old in existing_keys - set(counts.keys()):
            del var_store[old]

        for key in counts:
            if key not in var_store:
                var_store[key] = tk.BooleanVar(value=True)

            text_color = _Theme.TEXT_PRIMARY
            if color_map and key in color_map:
                text_color = color_map[key]

            cb = ctk.CTkCheckBox(
                self,
                text            = f"{key}  ({counts[key]})",
                variable        = var_store[key],
                font            = _Theme.FONT_BODY,
                text_color      = text_color,
                fg_color        = _Theme.ACCENT,
                hover_color     = _Theme.ACCENT_HOVER,
                checkmark_color = _Theme.BG_ROOT,
                border_color    = _Theme.BORDER,
                corner_radius   = 2,
                command         = self._on_change,
            )
            cb.grid(row=row, column=0, sticky="w",
                    padx=(_Theme.PAD_INNER, _Theme.PAD_SMALL), pady=2)
            row += 1
        return row


# ============================================================
# Individual flashcard widget
# ============================================================

class _CardWidget(ctk.CTkFrame):
    """
    Renders one _CardRecord as a card tile with VIEW and EDIT states.
    Calls on_change() after any save or delete.
    """

    _DIFF_ORDER = ["easy", "medium", "hard"]

    def __init__(
        self,
        parent,
        record:    _CardRecord,
        on_change: Callable[[], None],
        **kwargs,
    ) -> None:
        super().__init__(
            parent,
            fg_color      = _Theme.BG_CARD,
            corner_radius = _Theme.CORNER_RADIUS,
            border_width  = 1,
            border_color  = _Theme.BORDER,
            **kwargs,
        )
        self._record    = record
        self._on_change = on_change
        self._editing   = False
        self.grid_columnconfigure(0, weight=1)
        self._build_view()

    # â”€â”€ View mode â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

    def _build_view(self) -> None:
        self._clear()

        diff      = self._record.difficulty.lower()
        badge_col = _Theme.DIFF_COLORS.get(diff, _Theme.WARNING)

        header = ctk.CTkFrame(self, fg_color="transparent")
        header.grid(row=0, column=0, sticky="ew", padx=10, pady=(10, 4))
        header.grid_columnconfigure(1, weight=1)

        ctk.CTkLabel(
            header, text=diff.upper(),
            font=("Courier New", 9, "bold"), text_color=_Theme.BG_ROOT,
            fg_color=badge_col, corner_radius=2, width=48, height=18,
        ).grid(row=0, column=0, sticky="w")

        ctk.CTkLabel(
            header, text=self._record.source_name,
            font=_Theme.FONT_SMALL, text_color=_Theme.TEXT_SECONDARY, anchor="e",
        ).grid(row=0, column=1, sticky="e")

        ctk.CTkFrame(self, fg_color=_Theme.BORDER, height=1, corner_radius=0).grid(
            row=1, column=0, sticky="ew", padx=10)

        ctk.CTkLabel(
            self, text="Q", font=("Courier New", 10, "bold"),
            text_color=_Theme.ACCENT, anchor="w",
        ).grid(row=2, column=0, sticky="w", padx=(10, 0), pady=(8, 0))

        ctk.CTkLabel(
            self, text=self._wrap(self._record.question, 42),
            font=_Theme.FONT_CARD_Q, text_color=_Theme.TEXT_PRIMARY,
            anchor="w", justify="left", wraplength=_Theme.CARD_W - 30,
        ).grid(row=3, column=0, sticky="ew", padx=10, pady=(2, 4))

        ctk.CTkLabel(
            self, text="A", font=("Courier New", 10, "bold"),
            text_color=_Theme.TEXT_SECONDARY, anchor="w",
        ).grid(row=4, column=0, sticky="w", padx=(10, 0))

        ctk.CTkLabel(
            self, text=self._wrap(self._record.answer, 46),
            font=_Theme.FONT_CARD_A, text_color=_Theme.TEXT_MONO,
            anchor="w", justify="left", wraplength=_Theme.CARD_W - 30,
        ).grid(row=5, column=0, sticky="ew", padx=10, pady=(2, 8))

        if self._record.tags:
            tag_frame = ctk.CTkFrame(self, fg_color="transparent")
            tag_frame.grid(row=6, column=0, sticky="ew", padx=10, pady=(0, 6))
            for i, tag in enumerate(self._record.tags[:6]):
                ctk.CTkLabel(
                    tag_frame, text=f"#{tag}", font=_Theme.FONT_SMALL,
                    text_color=_Theme.TEXT_SECONDARY, fg_color=_Theme.BG_DROPZONE,
                    corner_radius=2, padx=5, pady=2,
                ).grid(row=0, column=i, padx=(0, 4), sticky="w")

        ctk.CTkFrame(self, fg_color=_Theme.BORDER, height=1, corner_radius=0).grid(
            row=7, column=0, sticky="ew", padx=10)

        action_row = ctk.CTkFrame(self, fg_color="transparent")
        action_row.grid(row=8, column=0, sticky="ew", padx=8, pady=(6, 8))
        action_row.grid_columnconfigure(0, weight=1)

        ctk.CTkButton(
            action_row, text="EDIT", font=("Courier New", 10, "bold"),
            width=60, height=26, corner_radius=_Theme.CORNER_RADIUS,
            fg_color="transparent", hover_color=_Theme.BG_DROPZONE_HV,
            text_color=_Theme.ACCENT, border_width=1, border_color=_Theme.BORDER,
            command=self._enter_edit,
        ).grid(row=0, column=0, sticky="w")

        ctk.CTkButton(
            action_row, text="X", font=("Trebuchet MS", 11, "bold"),
            width=28, height=26, corner_radius=_Theme.CORNER_RADIUS,
            fg_color="transparent", hover_color=_Theme.ERROR,
            text_color=_Theme.TEXT_SECONDARY, border_width=0,
            command=self._delete,
        ).grid(row=0, column=1, sticky="e")

    # â”€â”€ Edit mode â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

    def _build_edit(self) -> None:
        self._clear()
        self.configure(border_color=_Theme.ACCENT, fg_color=_Theme.BG_EDIT)

        ctk.CTkLabel(
            self, text="EDITING CARD", font=_Theme.FONT_MONO,
            text_color=_Theme.ACCENT, anchor="center",
        ).grid(row=0, column=0, sticky="ew", pady=(10, 6))

        ctk.CTkLabel(self, text="QUESTION", font=("Courier New", 9, "bold"),
                     text_color=_Theme.TEXT_SECONDARY, anchor="w",
                     ).grid(row=1, column=0, sticky="w", padx=10)

        self._q_box = ctk.CTkTextbox(
            self, height=80, font=_Theme.FONT_CARD_Q, fg_color=_Theme.BG_PANEL,
            text_color=_Theme.TEXT_PRIMARY, border_color=_Theme.BORDER_FOCUS,
            border_width=1, corner_radius=_Theme.CORNER_RADIUS, wrap="word",
        )
        self._q_box.grid(row=2, column=0, sticky="ew", padx=10, pady=(2, 8))
        self._q_box.insert("1.0", self._record.question)

        ctk.CTkLabel(self, text="ANSWER", font=("Courier New", 9, "bold"),
                     text_color=_Theme.TEXT_SECONDARY, anchor="w",
                     ).grid(row=3, column=0, sticky="w", padx=10)

        self._a_box = ctk.CTkTextbox(
            self, height=100, font=_Theme.FONT_CARD_A, fg_color=_Theme.BG_PANEL,
            text_color=_Theme.TEXT_PRIMARY, border_color=_Theme.BORDER_FOCUS,
            border_width=1, corner_radius=_Theme.CORNER_RADIUS, wrap="word",
        )
        self._a_box.grid(row=4, column=0, sticky="ew", padx=10, pady=(2, 8))
        self._a_box.insert("1.0", self._record.answer)

        ctk.CTkLabel(self, text="DIFFICULTY", font=("Courier New", 9, "bold"),
                     text_color=_Theme.TEXT_SECONDARY, anchor="w",
                     ).grid(row=5, column=0, sticky="w", padx=10)

        self._diff_var = tk.StringVar(value=self._record.difficulty)
        ctk.CTkOptionMenu(
            self, values=self._DIFF_ORDER, variable=self._diff_var,
            font=_Theme.FONT_BODY, fg_color=_Theme.BG_PANEL,
            button_color=_Theme.BG_CARD, button_hover_color=_Theme.BG_DROPZONE_HV,
            text_color=_Theme.TEXT_PRIMARY, dropdown_fg_color=_Theme.BG_PANEL,
            dropdown_text_color=_Theme.TEXT_PRIMARY,
            dropdown_hover_color=_Theme.BG_DROPZONE_HV,
            corner_radius=_Theme.CORNER_RADIUS, height=30, width=120,
        ).grid(row=6, column=0, sticky="w", padx=10, pady=(2, 8))

        ctk.CTkLabel(self, text="TAGS  (comma-separated)",
                     font=("Courier New", 9, "bold"),
                     text_color=_Theme.TEXT_SECONDARY, anchor="w",
                     ).grid(row=7, column=0, sticky="w", padx=10)

        self._tags_entry = ctk.CTkEntry(
            self, placeholder_text="e.g.  biology, cell-theory",
            font=_Theme.FONT_SMALL, fg_color=_Theme.BG_PANEL,
            text_color=_Theme.TEXT_PRIMARY, border_color=_Theme.BORDER,
            border_width=1, corner_radius=_Theme.CORNER_RADIUS, height=28,
        )
        self._tags_entry.grid(row=8, column=0, sticky="ew", padx=10, pady=(2, 10))
        self._tags_entry.insert(0, ", ".join(self._record.tags))

        btn_row = ctk.CTkFrame(self, fg_color="transparent")
        btn_row.grid(row=9, column=0, sticky="ew", padx=8, pady=(0, 10))
        btn_row.grid_columnconfigure(0, weight=1)

        ctk.CTkButton(
            btn_row, text="SAVE", font=("Courier New", 10, "bold"),
            width=70, height=28, corner_radius=_Theme.CORNER_RADIUS,
            fg_color=_Theme.ACCENT, hover_color=_Theme.ACCENT_HOVER,
            text_color=_Theme.BG_ROOT, command=self._save_edit,
        ).grid(row=0, column=0, sticky="w")

        ctk.CTkButton(
            btn_row, text="CANCEL", font=("Courier New", 10, "bold"),
            width=70, height=28, corner_radius=_Theme.CORNER_RADIUS,
            fg_color="transparent", hover_color=_Theme.BG_DROPZONE_HV,
            text_color=_Theme.TEXT_SECONDARY, border_width=1,
            border_color=_Theme.BORDER, command=self._cancel_edit,
        ).grid(row=0, column=1, padx=(6, 0), sticky="w")

    # â”€â”€ State transitions â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

    def _enter_edit(self) -> None:
        self._editing = True
        self._build_edit()

    def _save_edit(self) -> None:
        q = _sanitize_llm_text(self._q_box.get("1.0", "end"))
        a = _sanitize_llm_text(self._a_box.get("1.0", "end"))
        t = [_sanitize_tag_text(t) for t in self._tags_entry.get().split(",") if t.strip()]
        t = [tag for tag in t if tag]
        d = self._diff_var.get()

        if not q or not a:
            self.configure(border_color=_Theme.ERROR)
            self.after(800, lambda: self.configure(border_color=_Theme.ACCENT))
            return

        self._record.question   = q
        self._record.answer     = a
        self._record.tags       = t
        self._record.difficulty = d
        self._editing = False
        self.configure(border_color=_Theme.BORDER, fg_color=_Theme.BG_CARD)
        self._build_view()
        self._on_change()

    def _cancel_edit(self) -> None:
        self._editing = False
        self.configure(border_color=_Theme.BORDER, fg_color=_Theme.BG_CARD)
        self._build_view()

    def _delete(self) -> None:
        self._record.deleted = True
        self.grid_remove()
        self._on_change()

    # â”€â”€ Utility â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

    def _clear(self) -> None:
        for w in self.winfo_children():
            w.destroy()

    @staticmethod
    def _wrap(text: str, limit: int) -> str:
        return text if len(text) <= limit * 3 else text[: limit * 3] + "..."


# ============================================================
# Failed-chunk banner
# ============================================================

class _FailureBanner(ctk.CTkFrame):
    """Compact error notice for each failed GenerationResult."""

    def __init__(self, parent, result: "GenerationResult", **kwargs) -> None:
        super().__init__(
            parent, fg_color="#2a1a18", corner_radius=_Theme.CORNER_RADIUS,
            border_width=1, border_color=_Theme.ERROR, **kwargs,
        )
        src  = Path(result.source_file).name
        err  = result.error or "Unknown error"
        ctk.CTkLabel(
            self,
            text       = f"Chunk {result.chunk_index} of '{src}' failed - {err}",
            font       = _Theme.FONT_MONO,
            text_color = _Theme.ERROR,
            anchor     = "w",
            justify    = "left",
            wraplength = 760,
        ).pack(padx=12, pady=8, anchor="w")


# ============================================================
# Flashcard grid
# ============================================================

class _FlashcardGrid(ctk.CTkScrollableFrame):
    """Scrollable 2-column card grid with failure banners above."""

    def __init__(self, parent, on_change: Callable[[], None], **kwargs) -> None:
        super().__init__(
            parent, fg_color=_Theme.BG_ROOT, corner_radius=0,
            scrollbar_button_color=_Theme.BORDER,
            scrollbar_button_hover_color=_Theme.ACCENT,
            **kwargs,
        )
        self._on_change    = on_change
        self._records:      List[_CardRecord]       = []
        self._failed:       List["GenerationResult"] = []
        self._card_widgets: Dict[int, _CardWidget]   = {}
        self._visible_keys: Set[int]                 = set()
        self._content_row_offset                      = 0
        self._failed_widgets_built                    = False
        self._empty_label: Optional[ctk.CTkLabel]    = None
        self.grid_columnconfigure(0, weight=1)
        self.grid_columnconfigure(1, weight=1)

    def set_data(self, records: List[_CardRecord], failed: List["GenerationResult"]) -> None:
        self._records = records
        self._failed  = failed
        self._ensure_static_widgets()
        self._ensure_card_widgets()

    def _ensure_static_widgets(self) -> None:
        if self._failed_widgets_built:
            return

        current_row = 0
        if self._failed:
            ctk.CTkLabel(
                self, text="FAILED CHUNKS", font=_Theme.FONT_HEADING,
                text_color=_Theme.ERROR, anchor="w",
            ).grid(
                row=current_row, column=0, columnspan=2,
                sticky="ew", padx=_Theme.PAD_INNER, pady=(12, 4),
            )
            current_row += 1

            for result in self._failed:
                _FailureBanner(self, result).grid(
                    row=current_row, column=0, columnspan=2,
                    sticky="ew", padx=_Theme.PAD_INNER, pady=(0, _Theme.PAD_SMALL),
                )
                current_row += 1

            ctk.CTkFrame(self, fg_color=_Theme.BORDER, height=1, corner_radius=0).grid(
                row=current_row, column=0, columnspan=2,
                sticky="ew", padx=_Theme.PAD_INNER, pady=(4, 12),
            )
            current_row += 1

        self._content_row_offset = current_row
        self._empty_label = ctk.CTkLabel(
            self,
            text="No cards match the active filters.",
            font=_Theme.FONT_MONO,
            text_color=_Theme.TEXT_SECONDARY,
            anchor="center",
        )
        self._failed_widgets_built = True

    def _ensure_card_widgets(self) -> None:
        valid_keys = {id(record) for record in self._records}
        for key in list(self._card_widgets.keys()):
            if key not in valid_keys:
                self._card_widgets[key].destroy()
                del self._card_widgets[key]

        for record in self._records:
            key = id(record)
            if key not in self._card_widgets:
                self._card_widgets[key] = _CardWidget(self, record=record, on_change=self._on_change)

    def render(
        self,
        active_sources:      Set[str],
        active_difficulties: Set[str],
        active_tags:         Set[str],
    ) -> int:
        """Render filtered cards. Returns visible card count."""
        self._ensure_static_widgets()
        self._ensure_card_widgets()

        visible_records = [
            r for r in self._records
            if not r.deleted
            and r.source_name in active_sources
            and r.difficulty.lower() in active_difficulties
            and (not active_tags or bool(set(r.tags) & active_tags))
        ]
        visible_keys = {id(record) for record in visible_records}

        for key in self._visible_keys - visible_keys:
            widget = self._card_widgets.get(key)
            if widget is not None:
                widget.grid_remove()
        self._visible_keys = visible_keys

        if not visible_records:
            if self._empty_label is not None:
                self._empty_label.grid(
                    row=self._content_row_offset,
                    column=0,
                    columnspan=2,
                    padx=_Theme.PAD_OUTER,
                    pady=60,
                    sticky="ew",
                )
            return 0
        if self._empty_label is not None:
            self._empty_label.grid_remove()

        for index, record in enumerate(visible_records):
            col = index % 2
            row = self._content_row_offset + (index // 2)
            card = self._card_widgets[id(record)]
            card.grid(
                row    = row,
                column = col,
                sticky = "new",
                padx   = (
                    _Theme.PAD_INNER if col == 0 else _Theme.CARD_GAP // 2,
                    _Theme.CARD_GAP // 2 if col == 0 else _Theme.PAD_INNER,
                ),
                pady   = (0, _Theme.CARD_GAP),
            )

        return len(visible_records)


# ============================================================
# Export bar  â† UPDATED: format dropdown + unified confirm button
# ============================================================

class _ExportBar(ctk.CTkFrame):
    """
    Bottom action strip.

    Contains:
      â€¢ A CTkOptionMenu for selecting the export format
        (DOCX / PDF / PPTX / Anki / CSV / JSON â€” sourced from _FORMAT_LABELS)
      â€¢ A "Confirm & Save" button that fires on_confirm()
      â€¢ A live card-count label
    """

    def __init__(self, parent, on_confirm: Callable[[], None], **kwargs) -> None:
        super().__init__(parent, fg_color=_Theme.BG_PANEL, corner_radius=0, **kwargs)
        self._on_confirm = on_confirm
        self.grid_columnconfigure(2, weight=1)   # count label expands to fill
        self._build()

    def _build(self) -> None:
        # Top separator
        ctk.CTkFrame(self, fg_color=_Theme.BORDER, height=1, corner_radius=0).grid(
            row=0, column=0, columnspan=4, sticky="ew")

        # "FORMAT" label above the dropdown
        ctk.CTkLabel(
            self, text="FORMAT",
            font=("Courier New", 9, "bold"), text_color=_Theme.TEXT_SECONDARY, anchor="w",
        ).grid(row=1, column=0, padx=(_Theme.PAD_OUTER, 4), pady=(4, 0), sticky="sw")

        # â”€â”€ Format selector â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
        # The CTkOptionMenu reads _FORMAT_LABELS (derived from _EXPORT_FORMATS).
        # We store the selection in _format_var and read it in _confirm().
        self._format_var = tk.StringVar(value=_FORMAT_LABELS[0])

        self._format_menu = ctk.CTkOptionMenu(
            self,
            values               = _FORMAT_LABELS,
            variable             = self._format_var,
            font                 = ("Courier New", 11, "bold"),
            fg_color             = _Theme.BG_CARD,
            button_color         = _Theme.BG_DROPZONE,
            button_hover_color   = _Theme.BG_DROPZONE_HV,
            text_color           = _Theme.TEXT_PRIMARY,
            dropdown_fg_color    = _Theme.BG_PANEL,
            dropdown_text_color  = _Theme.TEXT_PRIMARY,
            dropdown_hover_color = _Theme.BG_DROPZONE_HV,
            corner_radius        = _Theme.CORNER_RADIUS,
            height               = _Theme.BTN_HEIGHT,
            width                = 190,
            # No command â€” we read the var at confirm time, not on change.
        )
        self._format_menu.grid(
            row=2, column=0, padx=(_Theme.PAD_OUTER, 8), pady=(2, 12), sticky="w")

        # â”€â”€ Confirm & Save button â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
        # A single button whose action is determined by the format dropdown.
        # Its label is temporarily changed to "âš  Error Saving" on failure.
        self._confirm_btn = ctk.CTkButton(
            self,
            text          = "CONFIRM & SAVE",
            font          = ("Courier New", 13, "bold"),
            height        = _Theme.BTN_HEIGHT,
            corner_radius = _Theme.CORNER_RADIUS,
            fg_color      = _Theme.ACCENT,
            hover_color   = _Theme.ACCENT_HOVER,
            text_color    = _Theme.BG_ROOT,
            width         = 190,
            command       = self._on_confirm,
        )
        self._confirm_btn.grid(
            row=2, column=1, padx=(0, _Theme.PAD_INNER), pady=(2, 12), sticky="w")

        # â”€â”€ Card count label â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
        self._count_lbl = ctk.CTkLabel(
            self, text="", font=_Theme.FONT_MONO,
            text_color=_Theme.TEXT_SECONDARY, anchor="e",
        )
        self._count_lbl.grid(
            row=2, column=2, sticky="e",
            padx=(_Theme.PAD_INNER, _Theme.PAD_OUTER))

    # â”€â”€ Public helpers called by VerificationHub â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

    @property
    def selected_format(self) -> str:
        """Return the currently selected format label string."""
        return self._format_var.get()

    def set_count(self, visible: int, total: int, deleted: int) -> None:
        self._count_lbl.configure(
            text=f"{visible} showing | {deleted} deleted | {total} total")

    def set_btn_text(self, text: str) -> None:
        """Temporarily change the confirm button label (e.g. on error)."""
        self._confirm_btn.configure(text=text)

    def reset_btn(self) -> None:
        """Restore the default confirm button label."""
        self._confirm_btn.configure(text="CONFIRM & SAVE")


# ============================================================
# Export engine  â† NEW: all file-writing logic
# ============================================================

class _ExportEngine:
    """Stateless file-export helpers for cards and summaries."""

    @staticmethod
    def _clean_card(card: _CardRecord) -> Dict[str, object]:
        tags = []
        seen = set()
        for tag in card.tags:
            cleaned_tag = _sanitize_tag_text(tag)
            if cleaned_tag and cleaned_tag not in seen:
                seen.add(cleaned_tag)
                tags.append(cleaned_tag)
        difficulty = str(card.difficulty or "medium").strip().lower()
        if difficulty not in {"easy", "medium", "hard"}:
            difficulty = "medium"
        return {
            "question": _sanitize_llm_text(card.question),
            "answer": _sanitize_llm_text(card.answer),
            "difficulty": difficulty,
            "tags": tags,
            "source_file": card.source_file,
            "source_name": card.source_name,
            "chunk_index": card.chunk_index,
        }

    @staticmethod
    def _clean_summary(summary: _SummaryRecord) -> Dict[str, object]:
        return {
            "source_file": summary.source_file,
            "source_name": summary.source_name,
            "chunk_index": summary.chunk_index,
            "summary": _sanitize_llm_text(summary.summary),
        }

    @staticmethod
    def _build_export_lines(cards: List[_CardRecord], summaries: List[_SummaryRecord]) -> List[str]:
        cards_clean = [_ExportEngine._clean_card(card) for card in cards]
        summaries_clean = [_ExportEngine._clean_summary(summary) for summary in summaries]

        lines: List[str] = []
        lines.append("MemoForge Export")
        lines.append(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
        lines.append("")

        if summaries_clean:
            lines.append("Summaries")
            lines.append("=========")
            for idx, summary in enumerate(summaries_clean, start=1):
                lines.append(f"{idx}. {summary['source_name']} (chunk {summary['chunk_index']})")
                lines.extend([part for part in str(summary["summary"]).splitlines() if part.strip()])
                lines.append("")

        if cards_clean:
            lines.append("Flashcards")
            lines.append("==========")
            for idx, card in enumerate(cards_clean, start=1):
                lines.append(f"{idx}. Q: {card['question']}")
                lines.append(f"   A: {card['answer']}")
                lines.append(f"   Difficulty: {card['difficulty']}")
                if card["tags"]:
                    lines.append(f"   Tags: {', '.join(card['tags'])}")
                lines.append(f"   Source: {card['source_name']} (chunk {card['chunk_index']})")
                lines.append("")

        if not cards and not summaries:
            lines.append("No generated content to export.")

        return lines

    @staticmethod
    def export_json(cards: List[_CardRecord], summaries: List[_SummaryRecord], path: str) -> None:
        cards_clean = [_ExportEngine._clean_card(card) for card in cards]
        summaries_clean = [_ExportEngine._clean_summary(summary) for summary in summaries]
        sources = sorted({str(r["source_file"]) for r in cards_clean} | {str(s["source_file"]) for s in summaries_clean})
        payload = {
            "memoforge_version": "4",
            "export_sources": sources,
            "total_cards": len(cards_clean),
            "total_summaries": len(summaries_clean),
            "summaries": [
                {
                    "source_file": s["source_file"],
                    "chunk_index": s["chunk_index"],
                    "summary": s["summary"],
                }
                for s in summaries_clean
            ],
            "cards": [
                {
                    "front": r["question"],
                    "back": r["answer"],
                    "difficulty": r["difficulty"],
                    "tags": r["tags"],
                    "source_file": r["source_file"],
                    "chunk_index": r["chunk_index"],
                }
                for r in cards_clean
            ],
        }
        with open(path, "w", encoding="utf-8") as fh:
            json.dump(payload, fh, indent=2, ensure_ascii=False)
        logger.info("JSON exported: %d cards, %d summaries -> %s", len(cards_clean), len(summaries_clean), path)

    @staticmethod
    def export_csv(cards: List[_CardRecord], summaries: List[_SummaryRecord], path: str) -> None:
        _ = summaries
        cards_clean = [_ExportEngine._clean_card(card) for card in cards]
        with open(path, "w", newline="", encoding="utf-8-sig") as fh:
            writer = csv.writer(fh, delimiter=",", quoting=csv.QUOTE_MINIMAL)
            writer.writerow(["Front", "Back", "Difficulty", "Tags", "Source", "Chunk"])
            for r in cards_clean:
                writer.writerow(
                    [r["question"], r["answer"], r["difficulty"], "; ".join(r["tags"]), r["source_file"], r["chunk_index"]]
                )
        logger.info("CSV exported: %d cards -> %s", len(cards_clean), path)

    @staticmethod
    def export_anki(cards: List[_CardRecord], summaries: List[_SummaryRecord], path: str) -> None:
        _ = summaries
        cards_clean = [_ExportEngine._clean_card(card) for card in cards]
        if not _GENANKI_AVAILABLE:
            raise ImportError(
                "The 'genanki' library is required for Anki export.\n"
                "Install it with:  pip install genanki"
            )

        model = genanki.Model(
            1_607_392_319,
            "MemoForge Basic",
            fields=[{"name": "Front"}, {"name": "Back"}],
            templates=[
                {
                    "name": "Card 1",
                    "qfmt": "{{Front}}",
                    "afmt": "{{FrontSide}}<hr id='answer'>{{Back}}",
                }
            ],
            css=".card { font-family: 'Trebuchet MS', sans-serif; font-size: 18px; text-align: left; }",
        )

        deck = genanki.Deck(random.randrange(1 << 30, 1 << 31), "MemoForge Export")
        for r in cards_clean:
            deck.add_note(
                genanki.Note(
                    model=model,
                    fields=[str(r["question"]), str(r["answer"])],
                    tags=[str(t).replace(" ", "_") for t in list(r["tags"])],
                )
            )
        genanki.Package(deck).write_to_file(path)
        logger.info("Anki .apkg exported: %d cards -> %s", len(cards_clean), path)

    @staticmethod
    def export_docx(cards: List[_CardRecord], summaries: List[_SummaryRecord], path: str) -> None:
        cards_clean = [_ExportEngine._clean_card(card) for card in cards]
        summaries_clean = [_ExportEngine._clean_summary(summary) for summary in summaries]
        try:
            from docx import Document
            from docx.shared import Pt
        except ImportError as exc:
            raise ImportError(
                "python-docx is required for DOCX export.\n"
                "Install it with: pip install python-docx"
            ) from exc

        document = Document()
        document.add_heading("MemoForge Study Notes", level=1)
        meta = document.add_paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
        if meta.runs:
            meta.runs[0].font.size = Pt(10)
        document.add_paragraph("")

        if summaries_clean:
            document.add_heading("Summaries", level=2)
            for summary in summaries_clean:
                document.add_heading(f"{summary['source_name']} - chunk {summary['chunk_index']}", level=3)
                summary_text = _sanitize_llm_text(str(summary["summary"]))
                sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+", summary_text) if s.strip()]
                if len(sentences) <= 1:
                    document.add_paragraph(summary_text)
                else:
                    intro = document.add_paragraph("Key points:")
                    if intro.runs:
                        intro.runs[0].bold = True
                    for sentence in sentences:
                        document.add_paragraph(sentence, style="List Bullet")
                document.add_paragraph("")

        if cards_clean:
            document.add_heading("Flashcards", level=2)
            for idx, card in enumerate(cards_clean, start=1):
                document.add_heading(f"Card {idx}", level=3)
                q_para = document.add_paragraph()
                q_para.add_run("Question: ").bold = True
                q_para.add_run(str(card["question"]))

                a_para = document.add_paragraph()
                a_para.add_run("Answer: ").bold = True
                a_para.add_run(str(card["answer"]))

                d_para = document.add_paragraph()
                d_para.add_run("Difficulty: ").bold = True
                d_para.add_run(str(card["difficulty"]).title())
                if card["tags"]:
                    t_para = document.add_paragraph()
                    t_para.add_run("Tags: ").bold = True
                    t_para.add_run(", ".join([f"#{tag}" for tag in list(card["tags"])]))
                s_para = document.add_paragraph()
                s_para.add_run("Source: ").bold = True
                s_para.add_run(f"{card['source_name']} (chunk {card['chunk_index']})")
                document.add_paragraph("")

        if not cards and not summaries:
            document.add_paragraph("No generated content to export.")

        document.save(path)
        logger.info("DOCX exported: %d cards, %d summaries -> %s", len(cards_clean), len(summaries_clean), path)

    @staticmethod
    def export_pptx(cards: List[_CardRecord], summaries: List[_SummaryRecord], path: str) -> None:
        cards_clean = [_ExportEngine._clean_card(card) for card in cards]
        summaries_clean = [_ExportEngine._clean_summary(summary) for summary in summaries]
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from pptx.util import Inches, Pt
        except ImportError as exc:
            raise ImportError(
                "python-pptx is required for PPTX export.\n"
                "Install it with: pip install python-pptx"
            ) from exc

        def _rgb(hex_color: str) -> RGBColor:
            value = hex_color.strip().lstrip("#")
            return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))

        def _truncate(text: str, max_chars: int) -> str:
            clean = _sanitize_llm_text(text)
            if len(clean) <= max_chars:
                return clean
            return clean[: max_chars - 3].rstrip() + "..."

        def _set_slide_bg(slide, hex_color: str) -> None:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = _rgb(hex_color)

        def _add_textbox(slide, left: float, top: float, width: float, height: float):
            return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))

        def _add_section_slide(title: str, subtitle: str) -> None:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            _set_slide_bg(slide, "#F8FAFF")

            ribbon = slide.shapes.add_shape(
                1, Inches(0), Inches(0), Inches(13.33), Inches(0.7)
            )
            ribbon.fill.solid()
            ribbon.fill.fore_color.rgb = _rgb("#1F4E9E")
            ribbon.line.fill.background()

            box = _add_textbox(slide, 0.6, 0.1, 10.8, 0.45)
            tf = box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = _rgb("#FFFFFF")

            box = _add_textbox(slide, 0.9, 1.2, 11.8, 1.0)
            tf = box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(20)
            p.font.color.rgb = _rgb("#243248")

        presentation = Presentation()
        presentation.slide_width = Inches(13.33)
        presentation.slide_height = Inches(7.5)

        # Title slide
        title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        _set_slide_bg(title_slide, "#0F1B2E")
        title_slide.shapes.title.text = "MemoForge Study Deck"
        title_slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)
        title_slide.shapes.title.text_frame.paragraphs[0].font.bold = True
        title_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = _rgb("#FFFFFF")
        if len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = (
                f"{len(cards_clean)} flashcards | {len(summaries_clean)} summaries\n"
                f"Generated {datetime.now().strftime('%Y-%m-%d %H:%M')}"
            )
            para = title_slide.placeholders[1].text_frame.paragraphs[0]
            para.font.size = Pt(19)
            para.font.color.rgb = _rgb("#C9D6EC")

        if summaries_clean:
            _add_section_slide(
                "Summaries",
                "High-level notes generated from your source material.",
            )

        for summary in summaries_clean:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            _set_slide_bg(slide, "#FFFFFF")

            top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.65))
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = _rgb("#1F4E9E")
            top_bar.line.fill.background()

            heading = _add_textbox(slide, 0.5, 0.1, 12.3, 0.45)
            tf = heading.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = f"Summary: {summary['source_name']} (chunk {summary['chunk_index']})"
            p.font.size = Pt(21)
            p.font.bold = True
            p.font.color.rgb = _rgb("#FFFFFF")

            body_shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.0), Inches(12.3), Inches(5.9))
            body_shape.fill.solid()
            body_shape.fill.fore_color.rgb = _rgb("#F4F7FC")
            body_shape.line.color.rgb = _rgb("#D4DFEF")
            body_shape.line.width = Pt(1.2)

            body_tf = body_shape.text_frame
            body_tf.clear()
            body_tf.word_wrap = True
            summary_text = _truncate(str(summary["summary"]), 1800)
            sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+", summary_text) if s.strip()]
            if not sentences:
                sentences = [summary_text]
            for idx, sentence in enumerate(sentences[:8]):
                para = body_tf.paragraphs[0] if idx == 0 else body_tf.add_paragraph()
                para.text = sentence
                para.level = 0
                para.font.size = Pt(20)
                para.font.color.rgb = _rgb("#243248")
                para.space_after = Pt(8)

        if cards_clean:
            _add_section_slide(
                "Flashcards",
                "Each slide contains one question-answer pair for quick revision.",
            )

        for idx, card in enumerate(cards_clean, start=1):
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            _set_slide_bg(slide, "#FFFFFF")

            header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.62))
            header.fill.solid()
            header.fill.fore_color.rgb = _rgb("#123A73")
            header.line.fill.background()

            header_text = _add_textbox(slide, 0.5, 0.08, 12.3, 0.4)
            tf = header_text.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = f"Flashcard {idx}"
            p.font.bold = True
            p.font.size = Pt(20)
            p.font.color.rgb = _rgb("#FFFFFF")

            question_box = slide.shapes.add_shape(1, Inches(0.6), Inches(0.95), Inches(12.1), Inches(2.35))
            question_box.fill.solid()
            question_box.fill.fore_color.rgb = _rgb("#EAF2FF")
            question_box.line.color.rgb = _rgb("#BFD3F3")
            question_box.line.width = Pt(1.0)

            q_tf = question_box.text_frame
            q_tf.clear()
            q_tf.word_wrap = True
            q_label = q_tf.paragraphs[0]
            q_label.text = "Question"
            q_label.font.bold = True
            q_label.font.size = Pt(17)
            q_label.font.color.rgb = _rgb("#19457D")

            q_text = q_tf.add_paragraph()
            q_text.text = _truncate(str(card["question"]), 340)
            q_text.font.size = Pt(25)
            q_text.font.bold = True
            q_text.font.color.rgb = _rgb("#12233A")
            q_text.space_before = Pt(6)

            answer_box = slide.shapes.add_shape(1, Inches(0.6), Inches(3.55), Inches(12.1), Inches(2.6))
            answer_box.fill.solid()
            answer_box.fill.fore_color.rgb = _rgb("#F6F8FC")
            answer_box.line.color.rgb = _rgb("#D6DEEB")
            answer_box.line.width = Pt(1.0)

            a_tf = answer_box.text_frame
            a_tf.clear()
            a_tf.word_wrap = True
            a_label = a_tf.paragraphs[0]
            a_label.text = "Answer"
            a_label.font.bold = True
            a_label.font.size = Pt(16)
            a_label.font.color.rgb = _rgb("#3C4F6A")

            a_text = a_tf.add_paragraph()
            a_text.text = _truncate(str(card["answer"]), 620)
            a_text.font.size = Pt(18)
            a_text.font.color.rgb = _rgb("#243248")
            a_text.space_before = Pt(5)

            footer = _add_textbox(slide, 0.65, 6.45, 12.0, 0.42)
            f_tf = footer.text_frame
            f_tf.clear()
            tags_text = ", ".join([_sanitize_tag_text(str(tag)) for tag in list(card["tags"]) if _sanitize_tag_text(str(tag))])
            info = (
                f"Difficulty: {str(card['difficulty']).title()} | Source: {card['source_name']}"
                f"{f' | Tags: {tags_text}' if tags_text else ''}"
            )
            p = f_tf.paragraphs[0]
            p.text = info
            p.font.size = Pt(12)
            p.font.color.rgb = _rgb("#4A5C75")
            p.alignment = PP_ALIGN.LEFT

            if card["tags"]:
                tag_bar = _add_textbox(slide, 0.65, 6.85, 12.0, 0.35)
                t_tf = tag_bar.text_frame
                t_tf.clear()
                p = t_tf.paragraphs[0]
                p.text = "Keywords: " + ", ".join(
                    [f"#{_sanitize_tag_text(str(tag))}" for tag in list(card["tags"]) if _sanitize_tag_text(str(tag))]
                )
                p.font.size = Pt(11)
                p.font.color.rgb = _rgb("#5E7088")

        presentation.save(path)
        logger.info("PPTX exported: %d cards, %d summaries -> %s", len(cards_clean), len(summaries_clean), path)

    @staticmethod
    def export_pdf(cards: List[_CardRecord], summaries: List[_SummaryRecord], path: str) -> None:
        try:
            import fitz  # type: ignore
        except ImportError as exc:
            raise ImportError(
                "PyMuPDF is required for PDF export.\n"
                "Install it with: pip install pymupdf"
            ) from exc

        import textwrap

        lines = _ExportEngine._build_export_lines(cards, summaries)
        doc = fitz.open()
        page = doc.new_page()
        x_margin = 48
        y = 60
        max_y = page.rect.height - 60
        line_height = 14

        for line in lines:
            wrapped_lines = textwrap.wrap(line, width=100) if line else [""]
            for wrapped in wrapped_lines:
                if y > max_y:
                    page = doc.new_page()
                    y = 60
                page.insert_text((x_margin, y), wrapped, fontsize=11, fontname="helv")
                y += line_height

        doc.save(path)
        doc.close()
        logger.info("PDF exported: %d cards, %d summaries -> %s", len(cards), len(summaries), path)
# Main Verification Hub window
# ============================================================

class VerificationHub(ctk.CTkToplevel):
    """
    Module 4 entry point.

    Instantiate after the pipeline completes and call grab_set() to make
    it modal:

        hub = VerificationHub(parent=app, results=results)
        hub.grab_set()
        hub.wait_window()

    Parameters
    ----------
    parent : ctk.CTk | ctk.CTkToplevel
        Parent window (the MemoForge main frame).
    results : List[GenerationResult]
        Output from VLMPipeline.run().
    on_close : Callable[[], None] | None
        Optional callback invoked when the window closes.
    """

    _WIN_W = 1160
    _WIN_H = 780
    _MIN_W = 900
    _MIN_H = 600

    def __init__(
        self,
        parent,
        results:  List["GenerationResult"],
        on_close: Optional[Callable[[], None]] = None,
        **kwargs,
    ) -> None:
        super().__init__(parent, **kwargs)

        self._results  = results
        self._on_close = on_close
        self._records: List[_CardRecord]       = []
        self._summaries: List[_SummaryRecord]  = []
        self._failed:  List["GenerationResult"] = []
        self._refresh_job: Optional[str] = None
        self._total_latency = sum(getattr(r, "latency_seconds", 0.0) for r in results)
        self._ingest_results(results)

        self._configure_window()
        self._build()
        self._populate_filter()
        self._refresh()

        self.protocol("WM_DELETE_WINDOW", self._close)

    # â”€â”€ Setup â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

    def _ingest_results(self, results: List["GenerationResult"]) -> None:
        for result in results:
            if not result.succeeded:
                self._failed.append(result)
                continue
            if result.summary and result.summary.strip():
                self._summaries.append(
                    _SummaryRecord(
                        source_file=result.source_file,
                        chunk_index=result.chunk_index,
                        summary=_sanitize_llm_text(result.summary),
                    )
                )
            for card in result.flashcards:
                self._records.append(
                    _CardRecord(
                        card        = card,
                        source_file = result.source_file,
                        chunk_index = result.chunk_index,
                    )
                )

    def _configure_window(self) -> None:
        self.title("MemoForge - Verification & Export Hub")
        self.geometry(f"{self._WIN_W}x{self._WIN_H}")
        self.minsize(self._MIN_W, self._MIN_H)
        self.configure(fg_color=_Theme.BG_ROOT)
        self.resizable(True, True)
        self._center_on_parent()

    def _center_on_parent(self) -> None:
        self.update_idletasks()
        try:
            px = self.master.winfo_x()
            py = self.master.winfo_y()
            pw = self.master.winfo_width()
            ph = self.master.winfo_height()
            x  = px + (pw - self._WIN_W) // 2
            y  = py + (ph - self._WIN_H) // 2
        except Exception:
            sw = self.winfo_screenwidth()
            sh = self.winfo_screenheight()
            x  = (sw - self._WIN_W) // 2
            y  = (sh - self._WIN_H) // 2
        self.geometry(f"{self._WIN_W}x{self._WIN_H}+{max(0, x)}+{max(0, y)}")

    # â”€â”€ Layout build â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

    def _build(self) -> None:
        self.grid_rowconfigure(2, weight=1)
        self.grid_columnconfigure(1, weight=1)

        self._build_title_bar()

        self._stats = _StatsBanner(self)
        self._stats.grid(row=1, column=0, columnspan=2, sticky="ew")

        self._filter = _FilterPanel(self, on_filter_change=self._schedule_refresh)
        self._filter.grid(row=2, column=0, sticky="nsew")

        # Vertical divider between filter sidebar and card grid
        ctk.CTkFrame(self, fg_color=_Theme.BORDER, width=1, corner_radius=0).grid(
            row=2, column=0, sticky="nse")

        self._grid = _FlashcardGrid(self, on_change=self._schedule_refresh)
        self._grid.grid(row=2, column=1, sticky="nsew")
        self._grid.set_data(self._records, self._failed)

        # Export bar: passes our _confirm handler so the bar's button fires it
        self._export_bar = _ExportBar(self, on_confirm=self._confirm)
        self._export_bar.grid(row=3, column=0, columnspan=2, sticky="ew")

    def _build_title_bar(self) -> None:
        bar = ctk.CTkFrame(
            self, fg_color=_Theme.BG_PANEL, corner_radius=0, height=52)
        bar.grid(row=0, column=0, columnspan=2, sticky="ew")
        bar.grid_columnconfigure(1, weight=1)
        bar.grid_propagate(False)

        ctk.CTkLabel(
            bar, text="MF", font=("Courier New", 18, "bold"), text_color=_Theme.ACCENT,
        ).grid(row=0, column=0, padx=(_Theme.PAD_OUTER, 6), pady=10, sticky="w")

        ctk.CTkLabel(
            bar, text="VERIFICATION & EXPORT HUB",
            font=_Theme.FONT_DISPLAY, text_color=_Theme.TEXT_PRIMARY, anchor="w",
        ).grid(row=0, column=1, sticky="w", pady=10)

        ctk.CTkButton(
            bar, text="SELECT ALL", font=("Courier New", 10, "bold"),
            width=100, height=28, corner_radius=_Theme.CORNER_RADIUS,
            fg_color="transparent", hover_color=_Theme.BG_DROPZONE_HV,
            text_color=_Theme.TEXT_SECONDARY, border_width=1,
            border_color=_Theme.BORDER, command=lambda: self._filter.select_all(),
        ).grid(row=0, column=2, padx=8, pady=10)

        ctk.CTkButton(
            bar, text="CLOSE", font=("Courier New", 10, "bold"),
            width=90, height=28, corner_radius=_Theme.CORNER_RADIUS,
            fg_color="transparent", hover_color=_Theme.ERROR,
            text_color=_Theme.TEXT_SECONDARY, border_width=1,
            border_color=_Theme.BORDER, command=self._close,
        ).grid(row=0, column=3, padx=(_Theme.PAD_SMALL, _Theme.PAD_OUTER), pady=10)

        ctk.CTkFrame(bar, fg_color=_Theme.BORDER, height=1, corner_radius=0).grid(
            row=1, column=0, columnspan=4, sticky="ew")

    # â”€â”€ Data & filter helpers â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

    def _populate_filter(self) -> None:
        self._filter.populate(self._records)

    def _schedule_refresh(self) -> None:
        if self._refresh_job is not None:
            try:
                self.after_cancel(self._refresh_job)
            except Exception:
                pass
            self._refresh_job = None
        self._refresh_job = self.after(33, self._refresh)

    def _refresh(self) -> None:
        """Recompute filter state and re-render the grid + stats banner."""
        self._refresh_job = None
        active_src  = self._filter.active_sources
        active_diff = self._filter.active_difficulties
        active_tags = self._filter.active_tags

        visible = self._grid.render(active_src, active_diff, active_tags)

        total   = len(self._records)
        deleted = sum(1 for r in self._records if r.deleted)
        sources = len({r.source_name for r in self._records if not r.deleted})
        failed  = len(self._failed)
        self._stats.refresh(
            total=total, visible=visible, deleted=deleted,
            sources=sources, failed_chunks=failed, total_latency=self._total_latency,
        )
        self._export_bar.set_count(visible, total, deleted)

    # â”€â”€ Confirm & export  â† THE FULLY IMPLEMENTED METHOD â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

    def _active_cards(self) -> List[_CardRecord]:
        """Return all non-deleted records regardless of the current filter."""
        return [r for r in self._records if not r.deleted]

    def _active_summaries(self) -> List[_SummaryRecord]:
        return list(self._summaries)

    def _confirm(self) -> None:
        """
        Unified export entry point â€” triggered by the "Confirm & Save" button.

        Flow
        ----
        1. Guard: abort with a warning if there is no generated content to export.
        2. Resolve the chosen format from the export bar's dropdown
           (reads _ExportBar.selected_format â†’ looks up _FORMAT_BY_LABEL).
        3. Open a native asksaveasfilename dialog.
           â€¢ defaultextension and filetypes are dynamically derived from
             the selected format dict â€” a single source of truth.
        4. Dispatch to the appropriate _ExportEngine static method via a
           label â†’ function lookup dict (extensible without if/elif chains).
        5. On success: show a non-blocking toast, then close the window
           after a short delay so the user can read the toast.
        6. On failure: flash "âš  Error Saving" on the button for 2.5 s,
           then show a messagebox with the specific error. The window
           stays open so the user can retry or pick a different path.
        """

        # â”€â”€ 1. Guard: nothing to export â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
        cards = self._active_cards()
        summaries = self._active_summaries()
        if not cards and not summaries:
            messagebox.showwarning(
                "Export - No Content",
                "There is no generated summary or flashcard content to export.",
                parent=self,
            )
            return

        # â”€â”€ 2. Resolve the chosen format â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
        # selected_format is a label string such as "Anki Deck (.apkg)".
        # _FORMAT_BY_LABEL maps it to the full format dict (ext + filetypes).
        fmt_label = self._export_bar.selected_format
        fmt       = _FORMAT_BY_LABEL.get(fmt_label)

        if fmt is None:
            # Should never happen: the menu values come from _FORMAT_LABELS.
            messagebox.showerror(
                "Export Error",
                f"Unknown format selected: '{fmt_label}'",
                parent=self,
            )
            return

        # â”€â”€ 3. Open the save-file dialog â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
        # Both defaultextension and filetypes are driven by the format dict,
        # so the dialog always presents the correct filter for the chosen format.
        path = filedialog.asksaveasfilename(
            parent           = self,
            title            = f"Export - {fmt_label}",
            defaultextension = fmt["ext"],
            filetypes        = fmt["filetypes"],
            initialfile      = f"memoforge_export{fmt['ext']}",
        )

        # User cancelled the dialog â€” leave the window open, do nothing.
        if not path:
            return

        # â”€â”€ 4. Dispatch to the correct exporter â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
        # Lookup dict: format label â†’ _ExportEngine static method.
        # Adding a new format = add an entry here + a method in _ExportEngine.
        dispatch: Dict[str, Callable[[List[_CardRecord], List[_SummaryRecord], str], None]] = {
            "DOCX (.docx)":         _ExportEngine.export_docx,
            "PDF (.pdf)":           _ExportEngine.export_pdf,
            "PowerPoint (.pptx)":   _ExportEngine.export_pptx,
            "Anki Deck (.apkg)":    _ExportEngine.export_anki,
            "CSV (.csv)":           _ExportEngine.export_csv,
            "JSON (.json)":         _ExportEngine.export_json,
        }

        export_fn = dispatch.get(fmt_label)
        if export_fn is None:
            messagebox.showerror(
                "Export Error",
                f"No exporter is implemented for '{fmt_label}'.",
                parent=self,
            )
            return

        # â”€â”€ 5/6. Execute â€” handle success and failure â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
        try:
            if fmt_label in {"Anki Deck (.apkg)", "CSV (.csv)"} and not cards:
                messagebox.showwarning(
                    "Export - No Flashcards",
                    f"'{fmt_label}' requires flashcards, but this run only has summaries.",
                    parent=self,
                )
                return

            export_fn(cards, summaries, path)

            # â”€â”€ Success â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
            # Show a transient toast, then close the hub after a short delay.
            self._show_success(
                f"{fmt_label} saved - {len(cards)} card(s) - {len(summaries)} summary item(s)\n{path}"
            )
            # 1.2 s delay lets the user read the toast before the window closes.
            self.after(1200, self._close)

        except ImportError as exc:
            # genanki not installed â€” give an actionable pip install message.
            logger.error("Export ImportError: %s", exc)
            self._flash_error_btn()
            messagebox.showerror("Missing Dependency", str(exc), parent=self)

        except OSError as exc:
            # Filesystem error: bad path, no permissions, disk full, etc.
            logger.error("Export OSError: %s", exc)
            self._flash_error_btn()
            messagebox.showerror(
                "Export Failed",
                f"Could not write file:\n{exc}",
                parent=self,
            )

        except Exception as exc:
            # Catch-all: unexpected failures (e.g. malformed card data, genanki bug).
            logger.exception("Unexpected export error: %s", exc)
            self._flash_error_btn()
            messagebox.showerror(
                "Unexpected Export Error",
                f"An unexpected error occurred:\n{type(exc).__name__}: {exc}\n\n"
                "Check the application log for details.",
                parent=self,
            )

    def _flash_error_btn(self) -> None:
        """
        Briefly relabel the Confirm & Save button to signal failure.
        Restores the original label after 2.5 seconds.

        This gives immediate, in-window visual feedback without closing
        the hub â€” the user can retry or change the save path.
        """
        self._export_bar.set_btn_text("ERROR SAVING")
        self.after(2500, self._export_bar.reset_btn)

    # â”€â”€ Utility â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

    def _show_success(self, message: str) -> None:
        """Non-blocking success toast rendered as a temporary overlay window."""
        toast = ctk.CTkToplevel(self)
        toast.overrideredirect(True)    # no title bar or window decorations
        toast.configure(fg_color=_Theme.BG_CARD)

        ctk.CTkLabel(
            toast,
            text          = f"Saved: {message}",
            font          = _Theme.FONT_MONO,
            text_color    = _Theme.SUCCESS,
            fg_color      = _Theme.BG_CARD,
            padx          = 20,
            pady          = 14,
            corner_radius = _Theme.CORNER_RADIUS,
        ).pack()

        # Anchor the toast near the bottom-left of the hub window
        self.update_idletasks()
        x = self.winfo_x() + 40
        y = self.winfo_y() + self.winfo_height() - 120
        toast.geometry(f"+{x}+{y}")

        # Auto-destroy after 2.8 s
        toast.after(2800, toast.destroy)

    def _close(self) -> None:
        if self._refresh_job is not None:
            try:
                self.after_cancel(self._refresh_job)
            except Exception:
                pass
            self._refresh_job = None
        if self._on_close:
            try:
                self._on_close()
            except Exception:
                pass
        try:
            self.grab_release()
        except Exception:
            pass
        self.destroy()


# ============================================================
# Convenience launcher (for standalone testing)
# ============================================================

def open_verification_hub(
    parent,
    results:  List["GenerationResult"],
    on_close: Optional[Callable[[], None]] = None,
    modal:    bool = True,
) -> VerificationHub:
    """
    Convenience wrapper for Module 5 integration.

    Usage:
        from verification_hub import open_verification_hub

        hub = open_verification_hub(
            parent   = self,
            results  = pipeline_results,
            on_close = lambda: self.set_generating_state(False),
            modal    = True,
        )
    """
    hub = VerificationHub(parent=parent, results=results, on_close=on_close)
    if modal:
        hub.grab_set()
    return hub


# ============================================================
# Standalone demo entry point
# ============================================================

def _make_demo_results() -> List["GenerationResult"]:
    """Synthetic data for smoke-testing without the ML stack."""

    topics = [
        ("What is mitosis?",
         "Mitosis is the process of cell division resulting in two identical daughter cells.",
         "easy", ["biology", "cell-theory"]),
        ("Define osmosis.",
         "Osmosis is the movement of water molecules through a semipermeable membrane from an area of low solute concentration to high.",
         "medium", ["biology", "chemistry"]),
        ("What is Newton's Second Law?",
         "F = ma. The net force on an object equals its mass times its acceleration.",
         "easy", ["physics", "mechanics"]),
        ("Explain the central dogma of molecular biology.",
         "DNA is transcribed to RNA, which is translated into protein.",
         "medium", ["biology", "genetics"]),
        ("What is the Heisenberg Uncertainty Principle?",
         "It is impossible to simultaneously know both the exact position and momentum of a particle.",
         "hard", ["physics", "quantum-mechanics"]),
        ("Define entropy.",
         "Entropy is a measure of disorder or randomness in a thermodynamic system.",
         "medium", ["physics", "thermodynamics"]),
        ("What is photosynthesis?",
         "The process by which plants convert light energy into chemical energy stored as glucose.",
         "easy", ["biology", "botany"]),
        ("Explain covalent bonding.",
         "A covalent bond is a chemical bond involving the sharing of electron pairs between atoms.",
         "medium", ["chemistry"]),
        ("What is the law of conservation of mass?",
         "Matter cannot be created or destroyed; the total mass of reactants equals total mass of products.",
         "easy", ["chemistry", "physics"]),
        ("Define a derivative in calculus.",
         "The derivative measures the rate of change of a function with respect to a variable.",
         "hard", ["mathematics", "calculus"]),
    ]

    r1_cards = [Flashcard(q, a, d, t, i) for i, (q, a, d, t) in enumerate(topics[:7])]
    r2_cards = [Flashcard(q, a, d, t, i) for i, (q, a, d, t) in enumerate(topics[7:])]

    return [
        GenerationResult(
            chunk_index=0, source_file="/home/user/Documents/biology_notes.pdf",
            flashcards=r1_cards, summary="Overview of cell biology fundamentals.",
            was_vision_chunk=False, error=None, latency_seconds=4.31,
        ),
        GenerationResult(
            chunk_index=1, source_file="/home/user/Documents/physics_lecture.pdf",
            flashcards=r2_cards,
            summary="Introduction to Newtonian mechanics and thermodynamics.",
            was_vision_chunk=True, error=None, latency_seconds=5.87,
        ),
        # Intentionally failed chunk â€” exercises the failure banner
        GenerationResult(
            chunk_index=2, source_file="/home/user/Documents/scanned_handout.pdf",
            flashcards=[], summary="", was_vision_chunk=True,
            error="VLM output parse failure: JSON schema mismatch after retry.",
            latency_seconds=2.10,
        ),
    ]


if __name__ == "__main__":
    import sys
    logging.basicConfig(
        level=logging.DEBUG,
        format="%(asctime)s [%(levelname)s] %(name)s - %(message)s",
        stream=sys.stdout,
    )

    ctk.set_appearance_mode("dark")
    ctk.set_default_color_theme("dark-blue")

    root = ctk.CTk()
    root.title("MemoForge - Dev Shell (Module 4 test)")
    root.geometry("900x500")
    root.configure(fg_color=_Theme.BG_ROOT)

    demo_results = _make_demo_results()

    ctk.CTkLabel(
        root,
        text="Module 4 standalone demo.\nClick the button to open the Verification & Export Hub.",
        font=("Trebuchet MS", 13), text_color="#9a9585",
    ).pack(expand=True)

    ctk.CTkButton(
        root,
        text="OPEN VERIFICATION HUB",
        font=("Courier New", 13, "bold"),
        fg_color=_Theme.ACCENT, hover_color=_Theme.ACCENT_HOVER,
        text_color=_Theme.BG_ROOT, height=44,
        command=lambda: open_verification_hub(
            root, demo_results,
            on_close=lambda: logger.info("Hub closed."),
        ),
    ).pack(pady=20)

    root.mainloop()

